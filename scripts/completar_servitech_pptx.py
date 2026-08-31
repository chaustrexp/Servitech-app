import os
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def update_entire_presentation():
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    pptx_path = os.path.join(base_dir, "Servitech-app.pptx")

    prs = pptx.Presentation(pptx_path)

    # ── Paleta de Colores ────────────────────────────────────────────────────
    C_DARK = RGBColor(22, 22, 30)           # Texto oscuro (#16161E)
    C_WHITE = RGBColor(255, 255, 255)       # Blanco puro
    C_RED_CORAL = RGBColor(255, 128, 128)   # Coral / Rojo suave (Prioridad 25 / Crítico)
    C_ORANGE = RGBColor(255, 194, 102)      # Naranja cálido (Prioridad 16-20)
    C_YELLOW = RGBColor(255, 242, 163)      # Amarillo suave (Prioridad 9-12)
    C_SENA_GREEN = RGBColor(57, 169, 0)     # Verde SENA (#39A900)
    C_SENA_NAVY = RGBColor(0, 50, 77)       # Azul Marino SENA (#00324D)
    C_GREEN_PASS = RGBColor(46, 125, 50)    # Verde Aprobado Pruebas (#2E7D32)
    C_LIGHT_BG = RGBColor(248, 250, 252)    # Fila alternada (#F8FAFC)
    C_BORDER = RGBColor(203, 213, 225)      # Borde sutil (#CBD5E1)

    def set_slide_title(slide, title_text):
        title_shape = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape.top.inches < 1.6 and shape.left.inches < 2.0:
                title_shape = shape
                break

        if not title_shape:
            title_shape = slide.shapes.add_textbox(Inches(0.92), Inches(0.40), Inches(11.5), Inches(1.45))

        tf = title_shape.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Arial"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = C_DARK

    def clean_slide_body(slide):
        """Elimina formas del cuerpo conservando el título y el logo superior derecho."""
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.left.inches >= 11.5 and shape.top.inches <= 1.0:
                continue
            if shape.has_text_frame and shape.top.inches < 1.4 and shape.left.inches < 2.0:
                continue
            shapes_to_remove.append(shape)

        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)

    def set_speaker_notes(slide, text):
        slide.notes_slide.notes_text_frame.text = text

    # =========================================================================
    # SECCIÓN 1: FORMULACIÓN DEL PROYECTO (SLIDES 4 A 10)
    # =========================================================================

    # SLIDE 04: Planteamiento del Problema
    s04 = prs.slides[3]
    set_slide_title(s04, "Planteamiento del Problema")
    clean_slide_body(s04)
    box04 = s04.shapes.add_textbox(Inches(0.92), Inches(1.85), Inches(11.50), Inches(5.00))
    tf04 = box04.text_frame
    tf04.word_wrap = True
    tf04.margin_left = tf04.margin_top = tf04.margin_right = tf04.margin_bottom = 0
    tf04.clear()

    p = tf04.paragraphs[0]
    p.text = "En los talleres de servicio técnico no existe un sistema centralizado para gestionar turnos, citas y órdenes de reparación, lo que provoca:"
    p.font.name = "Arial"
    p.font.size = Pt(15.5)
    p.font.bold = False
    p.font.color.rgb = C_DARK
    p.space_after = Pt(12)

    problemas = [
        "• Gestión dispersa en cuadernos manuales, WhatsApp y hojas de cálculo no sincronizadas.",
        "• Largos tiempos de espera presenciales e incertidumbre en la fecha de entrega para los clientes.",
        "• Falta de trazabilidad en el estado del equipo dentro de los boxes técnicos (recepción, diagnóstico, reparación, entrega).",
        "• Mala coordinación operativa, tiempos muertos en mesas de trabajo y retrasos no notificados al cliente.",
        "• Descontrol en el inventario de repuestos utilizados y mermas operativas en el taller."
    ]
    for prob in problemas:
        p_prob = tf04.add_paragraph()
        p_prob.text = prob
        p_prob.font.name = "Arial"
        p_prob.font.size = Pt(14)
        p_prob.font.color.rgb = C_DARK
        p_prob.space_after = Pt(6)

    p_sol = tf04.add_paragraph()
    p_sol.text = "Solución: Desarrollar e implementar una plataforma web integral que centralice el agendamiento 24/7, la trazabilidad de turnos en boxes y el control de inventario en tiempo real."
    p_sol.font.name = "Arial"
    p_sol.font.size = Pt(14.5)
    p_sol.font.bold = True
    p_sol.font.color.rgb = C_DARK
    p_sol.space_before = Pt(10)
    set_speaker_notes(s04, "El problema central radica en la falta de un sistema digital de turnos y trazabilidad, generando fricción con los clientes y pérdidas de productividad.")

    # SLIDE 05: Justificación (Necesidad y Solución)
    s05 = prs.slides[4]
    set_slide_title(s05, "Justificación")
    clean_slide_body(s05)
    
    # Subtítulo introductorio
    box_sub05 = s05.shapes.add_textbox(Inches(0.92), Inches(1.80), Inches(11.50), Inches(0.45))
    p_sub05 = box_sub05.text_frame.paragraphs[0]
    p_sub05.text = "Se propone implementar una plataforma de software integral que consolide la atención al cliente y la operación técnica:"
    p_sub05.font.name = "Arial"
    p_sub05.font.size = Pt(14.5)
    p_sub05.font.color.rgb = C_DARK

    # Tabla Necesidad vs Solución
    headers05 = ["Necesidad Identificada", "Solución Implementada"]
    rows05 = [
        ["Información dispersa en registros físicos", "Unificación en base de datos relacional centralizada"],
        ["Incertidumbre en tiempos de entrega y SLA", "Cálculo automático de duración (45-90 min) y seguimiento en vivo"],
        ["Retrasos imprevistos e inasistencias de clientes", "Botón de aviso 'Llegaré tarde' (+15 min) y reasignación a walk-ins"],
        ["Tiempos muertos y desorganización de boxes", "Panel operativo con estados en tiempo real (Disponible, En Servicio, Pausa)"],
        ["Descontrol de inventario y mermas de repuestos", "Kardex automatizado de repuestos asociado a cada orden"],
        ["Riesgo de manipulación y alteración de datos", "Auditoría transaccional automática inmutable con triggers PostgreSQL"]
    ]
    t_shape05 = s05.shapes.add_table(len(rows05) + 1, 2, Inches(0.92), Inches(2.40), Inches(11.50), Inches(4.50))
    t05 = t_shape05.table
    t05.columns[0].width = Inches(5.20)
    t05.columns[1].width = Inches(6.30)
    for c_idx, h in enumerate(headers05):
        c = t05.cell(0, c_idx)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.fill.solid()
        c.fill.fore_color.rgb = C_WHITE
        p = c.text_frame.paragraphs[0]
        p.text = h
        p.font.name = "Arial"
        p.font.size = Pt(12.5)
        p.font.bold = True
        p.font.color.rgb = C_DARK
        p.alignment = PP_ALIGN.CENTER

    for r_idx, (nec, sol) in enumerate(rows05):
        bg = C_WHITE if r_idx % 2 == 0 else C_LIGHT_BG
        for c_idx, val in enumerate([nec, sol]):
            c = t05.cell(r_idx + 1, c_idx)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            c.fill.fore_color.rgb = bg
            p = c.text_frame.paragraphs[0]
            p.text = val
            p.font.name = "Arial"
            p.font.size = Pt(11.5)
            p.font.color.rgb = C_DARK
            if c_idx == 0:
                p.font.bold = True
    set_speaker_notes(s05, "La justificación se basa en contrastar cada dolor operativo del taller con una solución tecnológica concreta.")

    # SLIDE 06: Objetivo General
    s06 = prs.slides[5]
    set_slide_title(s06, "Objetivo General")
    clean_slide_body(s06)
    box06 = s06.shapes.add_textbox(Inches(1.20), Inches(2.30), Inches(10.80), Inches(4.00))
    tf06 = box06.text_frame
    tf06.word_wrap = True
    p06 = tf06.paragraphs[0]
    p06.text = "Desarrollar e implementar un sistema web centralizado para la gestión integral de citas, turnos de atención técnica y control de reparaciones en talleres de servicio técnico, mejorando la eficiencia operativa, el cumplimiento de los tiempos de entrega (SLA) y la trazabilidad del cliente."
    p06.font.name = "Arial"
    p06.font.size = Pt(21)
    p06.font.bold = False
    p06.font.color.rgb = C_DARK
    set_speaker_notes(s06, "El objetivo general abarca el desarrollo del software enfocado en la eficiencia operativa, SLA y satisfacción del cliente.")

    # SLIDE 07: Objetivos Específicos (Del Proyecto y Del Sistema)
    s07 = prs.slides[6]
    set_slide_title(s07, "Objetivos Específicos")
    clean_slide_body(s07)

    # Columna 1: Del Proyecto
    box_p = s07.shapes.add_textbox(Inches(0.92), Inches(1.85), Inches(5.60), Inches(5.00))
    tf_p = box_p.text_frame
    tf_p.word_wrap = True
    tf_p.clear()
    p_p0 = tf_p.paragraphs[0]
    p_p0.text = "Del proyecto:"
    p_p0.font.name = "Arial"
    p_p0.font.size = Pt(18)
    p_p0.font.bold = True
    p_p0.font.color.rgb = C_DARK
    p_p0.space_after = Pt(10)

    obj_proy = [
        "• Optimizar la recepción y asignación de equipos según disponibilidad técnica.",
        "• Reducir los tiempos muertos en los boxes de trabajo del taller.",
        "• Eliminar la incertidumbre de tiempos de entrega mediante acuerdos SLA.",
        "• Controlar el inventario de repuestos y reducir mermas operativas.",
        "• Mejorar la satisfacción del cliente eliminando filas presenciales."
    ]
    for op in obj_proy:
        p = tf_p.add_paragraph()
        p.text = op
        p.font.name = "Arial"
        p.font.size = Pt(13.5)
        p.font.color.rgb = C_DARK
        p.space_after = Pt(6)

    # Columna 2: Del Sistema
    box_s = s07.shapes.add_textbox(Inches(6.82), Inches(1.85), Inches(5.60), Inches(5.00))
    tf_s = box_s.text_frame
    tf_s.word_wrap = True
    tf_s.clear()
    p_s0 = tf_s.paragraphs[0]
    p_s0.text = "Del Sistema:"
    p_s0.font.name = "Arial"
    p_s0.font.size = Pt(18)
    p_s0.font.bold = True
    p_s0.font.color.rgb = C_DARK
    p_s0.space_after = Pt(10)

    obj_sys = [
        "• Desarrollar un módulo de agendamiento 24/7 con validación en tiempo real.",
        "• Construir el panel operativo de boxes con control de turnos y estados.",
        "• Implementar el enlace de contingencia 'Llegaré tarde' (+15 min).",
        "• Integrar el módulo de catálogo de servicios con tiempos SLA configurables.",
        "• Configurar la auditoría transaccional inmutable mediante triggers PostgreSQL.",
        "• Diseñar una interfaz intuitiva, accesible y multidispositivo."
    ]
    for osys in obj_sys:
        p = tf_s.add_paragraph()
        p.text = osys
        p.font.name = "Arial"
        p.font.size = Pt(13.5)
        p.font.color.rgb = C_DARK
        p.space_after = Pt(5)

    set_speaker_notes(s07, "Se dividen los objetivos entre los impactos de negocio del proyecto y los requerimientos funcionales del sistema.")

    # SLIDE 08: Alcance
    s08 = prs.slides[7]
    set_slide_title(s08, "Alcance")
    clean_slide_body(s08)
    box08 = s08.shapes.add_textbox(Inches(0.92), Inches(1.85), Inches(11.50), Inches(5.00))
    tf08 = box08.text_frame
    tf08.word_wrap = True
    tf08.clear()

    alcances = [
        "• Registro y autenticación de usuarios con control de roles (Administrador, Jefe de Taller, Técnico N1-N3, Cliente).",
        "• Motor de agendamiento 24/7 con selección de tipo de dispositivo, servicio requerido, técnico y franja horaria.",
        "• Panel de control técnico en tiempo real para asignación de turnos, inicio de servicio, pausas e incidencias en boxes.",
        "• Sistema de contingencias con notificación de retrasos ('Llegaré tarde') y reasignación automática para clientes presenciales (walk-ins).",
        "• Gestión de catálogo de servicios con parametrización de costos base, duraciones estimadas (45-90 min) y buffer SLA.",
        "• Módulo de inventario y Kardex de repuestos vinculado a las órdenes de trabajo ejecutadas.",
        "• Módulo de auditoría transaccional inmutable en PostgreSQL con registro de operaciones e historial JSONB.",
        "• Generación de reportes analíticos de productividad técnica y cumplimiento de SLA exportables en Excel y PDF."
    ]
    for idx, alc in enumerate(alcances):
        p = tf08.add_paragraph() if idx > 0 else tf08.paragraphs[0]
        p.text = alc
        p.font.name = "Arial"
        p.font.size = Pt(13.5)
        p.font.color.rgb = C_DARK
        p.space_after = Pt(6)
    set_speaker_notes(s08, "El alcance define el límite funcional del aplicativo en sus módulos de clientes, técnico, administrativo y auditoría.")

    # SLIDE 09: Riesgos
    s09 = prs.slides[8]
    set_slide_title(s09, "Riesgos")
    clean_slide_body(s09)
    box09 = s09.shapes.add_textbox(Inches(0.92), Inches(1.85), Inches(11.50), Inches(5.00))
    tf09 = box09.text_frame
    tf09.word_wrap = True
    tf09.clear()

    riesgos = [
        "• Resistencia al cambio por parte de clientes o personal técnico acostumbrados a la gestión en papel.",
        "• Interrupción o caídas temporales de la conectividad a internet en el establecimiento del taller.",
        "• Daño, obsolescencia o incompatibilidad de hardware en las estaciones de boxes de trabajo.",
        "• Desabastecimiento de repuestos críticos por demoras de proveedores externos.",
        "• Omisión operativa de los técnicos al no actualizar oportunamente los estados de las órdenes.",
        "• Sobrecarga de concurrencia y peticiones simultáneas en horas pico de reservas.",
        "• Incumplimiento de hitos en el cronograma de desarrollo y pruebas de integración."
    ]
    for idx, rsg in enumerate(riesgos):
        p = tf09.add_paragraph() if idx > 0 else tf09.paragraphs[0]
        p.text = rsg
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.color.rgb = C_DARK
        p.space_after = Pt(8)
    set_speaker_notes(s09, "Se identificaron los riesgos operativos, tecnológicos y humanos junto con sus planes de mitigación.")

    # SLIDE 10: Restricciones
    s10 = prs.slides[9]
    set_slide_title(s10, "Restricciones")
    clean_slide_body(s10)
    box10 = s10.shapes.add_textbox(Inches(0.92), Inches(1.85), Inches(11.50), Inches(5.00))
    tf10 = box10.text_frame
    tf10.word_wrap = True
    tf10.clear()

    restricciones = [
        "• Plazo estricto de desarrollo ajustado al cronograma académico del proyecto formativo SENA.",
        "• Stack tecnológico obligatorio: Python 3.12, Django 5.x, PostgreSQL 16 y TailwindCSS.",
        "• Requisito de ejecución en navegadores web modernos (Chrome, Firefox, Edge) sin instalación de software cliente.",
        "• Cumplimiento obligatorio de la Ley 1581 de 2012 de Protección de Datos Personales (Habeas Data).",
        "• Arquitectura optimizada para tiempos de respuesta web inferiores a 200 ms y sobrecarga de auditoría menor a 5 ms.",
        "• Acceso restringido por roles y permisos diferenciados (Administrador, Jefe de Taller, Técnico, Cliente)."
    ]
    for idx, rst in enumerate(restricciones):
        p = tf10.add_paragraph() if idx > 0 else tf10.paragraphs[0]
        p.text = rst
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.color.rgb = C_DARK
        p.space_after = Pt(8)
    set_speaker_notes(s10, "Las restricciones delimitan el marco tecnológico, legal, de rendimiento y temporal en el que opera el proyecto.")

    # =========================================================================
    # SECCIÓN 3: DISEÑO DE BASE DE DATOS (SLIDES 28, 29, 30)
    # =========================================================================

    # SLIDE 28: Diseño de la Base de Datos (Modelo Relacional / Diagrama)
    s28 = prs.slides[27]
    set_slide_title(s28, "Diseño de la Base de Datos")
    clean_slide_body(s28)
    
    # Subtítulo explicativo del Modelo Relacional
    sub28 = s28.shapes.add_textbox(Inches(0.92), Inches(1.80), Inches(11.50), Inches(0.45))
    p_sub28 = sub28.text_frame.paragraphs[0]
    p_sub28.text = "Estructura del Modelo Entidad-Relación y Tablas Transaccionales de ServiTech:"
    p_sub28.font.name = "Arial"
    p_sub28.font.size = Pt(15.5)
    p_sub28.font.bold = True
    p_sub28.font.color.rgb = C_DARK

    headers28 = ["Entidad / Tabla", "Tipo", "Relación Principal", "Propósito en el Negocio"]
    rows28 = [
        ["turnos_usuario", "Maestra", "1:N con turnos_cita", "Gestión de roles: Admin, Jefe de Taller, Técnico N1-N3 y Cliente"],
        ["turnos_servicio", "Catálogo", "1:N con turnos_cita", "Catálogo de servicios con costos base y duraciones SLA (45-90 min)"],
        ["turnos_cita", "Transaccional", "N:1 usuario, N:1 servicio", "Registro de citas, asignación de box/técnico y estados de atención"],
        ["turnos_repuesto", "Inventario", "1:N con cita_repuesto", "Kardex de repuestos, stock disponible, precio unitario y compatibilidad"],
        ["turnos_citarepuesto", "Intermedia", "N:1 cita, N:1 repuesto", "Detalle de piezas y repuestos instalados por orden técnica"],
        ["turnos_box", "Operativa", "1:1 con técnico asignado", "Control de boxes físicos de trabajo y estado de ocupación"],
        ["auditoria_log", "Auditoría", "Triggers nativos PL/pgSQL", "Registro inmutable con snapshot JSONB de toda mutación (INSERT/UPDATE/DELETE)"]
    ]
    t_shape28 = s28.shapes.add_table(len(rows28) + 1, 4, Inches(0.92), Inches(2.40), Inches(11.50), Inches(4.50))
    t28 = t_shape28.table
    t28.columns[0].width = Inches(2.40)
    t28.columns[1].width = Inches(1.80)
    t28.columns[2].width = Inches(3.10)
    t28.columns[3].width = Inches(4.20)
    for c_idx, h in enumerate(headers28):
        c = t28.cell(0, c_idx)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.fill.solid()
        c.fill.fore_color.rgb = C_WHITE
        p = c.text_frame.paragraphs[0]
        p.text = h
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = C_DARK
        p.alignment = PP_ALIGN.CENTER

    for r_idx, r_data in enumerate(rows28):
        bg = C_WHITE if r_idx % 2 == 0 else C_LIGHT_BG
        for c_idx, val in enumerate(r_data):
            c = t28.cell(r_idx + 1, c_idx)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            c.fill.fore_color.rgb = bg
            p = c.text_frame.paragraphs[0]
            p.text = val
            p.font.name = "Arial"
            p.font.size = Pt(11)
            p.font.color.rgb = C_DARK
            if c_idx in [0, 1]:
                p.font.bold = True
                if c_idx == 1:
                    p.alignment = PP_ALIGN.CENTER
    set_speaker_notes(s28, "Diseño de base de datos normalizado en tercera forma normal (3FN) optimizado para PostgreSQL.")

    # SLIDE 29: Diseño de la Base de Datos - Tabla `turnos_cita`
    s29 = prs.slides[28]
    set_slide_title(s29, "Diseño de la Base de Datos – Tabla")
    clean_slide_body(s29)

    sub29 = s29.shapes.add_textbox(Inches(0.92), Inches(1.80), Inches(11.50), Inches(0.45))
    p_sub29 = sub29.text_frame.paragraphs[0]
    p_sub29.text = "Diccionario de Datos: Tabla Principal `turnos_cita`"
    p_sub29.font.name = "Arial"
    p_sub29.font.size = Pt(16)
    p_sub29.font.bold = True
    p_sub29.font.color.rgb = C_DARK

    headers_dict = ["Nombre del Campo", "Tipo de Dato", "Longitud", "Llave", "Descripción"]
    rows29 = [
        ["id", "BIGINT (Auto)", "-", "PK", "Identificador único secuencial de la cita / orden"],
        ["cliente_id", "BIGINT", "-", "FK", "Referencia al usuario cliente solicitante (turnos_usuario)"],
        ["tecnico_id", "BIGINT", "-", "FK (Null)", "Técnico asignado a la orden técnica o box"],
        ["servicio_id", "BIGINT", "-", "FK", "Servicio técnico contratado con duración SLA base"],
        ["dispositivo_tipo", "VARCHAR", "50", "-", "Tipo de equipo: Celular, Laptop, PC, Consola"],
        ["marca_modelo", "VARCHAR", "150", "-", "Marca y modelo comercial del equipo a reparar"],
        ["serial_imei", "VARCHAR", "100", "-", "Número de serie o código IMEI para trazabilidad"],
        ["fecha_cita", "DATE", "-", "-", "Fecha programada para la atención presencial"],
        ["hora_inicio", "TIME", "-", "-", "Hora de inicio de la franja horaria agendada"],
        ["estado", "VARCHAR", "30", "-", "Confirmada, En Box, Pausada, Retrasada, Finalizada"]
    ]
    t_shape29 = s29.shapes.add_table(len(rows29) + 1, 5, Inches(0.92), Inches(2.35), Inches(11.50), Inches(4.55))
    t29 = t_shape29.table
    t29.columns[0].width = Inches(2.40)
    t29.columns[1].width = Inches(1.80)
    t29.columns[2].width = Inches(1.20)
    t29.columns[3].width = Inches(1.10)
    t29.columns[4].width = Inches(5.00)
    for c_idx, h in enumerate(headers_dict):
        c = t29.cell(0, c_idx)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.fill.solid()
        c.fill.fore_color.rgb = C_WHITE
        p = c.text_frame.paragraphs[0]
        p.text = h
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = C_DARK
        p.alignment = PP_ALIGN.CENTER

    for r_idx, r_data in enumerate(rows29):
        bg = C_WHITE if r_idx % 2 == 0 else C_LIGHT_BG
        for c_idx, val in enumerate(r_data):
            c = t29.cell(r_idx + 1, c_idx)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            c.fill.fore_color.rgb = bg
            p = c.text_frame.paragraphs[0]
            p.text = val
            p.font.name = "Arial"
            p.font.size = Pt(10.5)
            p.font.color.rgb = C_DARK
            if c_idx == 0:
                p.font.bold = True
            elif c_idx in [1, 2, 3]:
                p.alignment = PP_ALIGN.CENTER
                if c_idx == 3 and "PK" in val:
                    p.font.bold = True
    set_speaker_notes(s29, "Diccionario de datos de la tabla transaccional central que almacena todas las reservas y órdenes del taller.")

    # SLIDE 30: Diseño de la Base de Datos - Tablas `turnos_servicio` y `auditoria_log`
    s30 = prs.slides[29]
    set_slide_title(s30, "Diseño de la Base de Datos – Tabla")
    clean_slide_body(s30)

    sub30 = s30.shapes.add_textbox(Inches(0.92), Inches(1.80), Inches(11.50), Inches(0.45))
    p_sub30 = sub30.text_frame.paragraphs[0]
    p_sub30.text = "Diccionario de Datos: `turnos_servicio` (Catálogo SLA) y `auditoria_log`"
    p_sub30.font.name = "Arial"
    p_sub30.font.size = Pt(16)
    p_sub30.font.bold = True
    p_sub30.font.color.rgb = C_DARK

    rows30 = [
        ["[turnos_servicio] id", "BIGINT", "-", "PK", "Identificador único del servicio de catálogo"],
        ["nombre", "VARCHAR", "100", "-", "Nombre del servicio técnico (ej. Cambio de Pantalla)"],
        ["costo_base", "DECIMAL", "10,2", "-", "Tarifa base de mano de obra en pesos"],
        ["duracion_estimada_min", "INTEGER", "-", "-", "Tiempo SLA estimado de ejecución técnica (45-90 min)"],
        ["buffer_minutos", "INTEGER", "-", "-", "Margen de tolerancia antes de la siguiente cita (15 min)"],
        ["[auditoria_log] id", "BIGINT", "-", "PK", "Identificador secuencial del evento de auditoría"],
        ["tabla_afectada", "VARCHAR", "60", "-", "Nombre de la tabla mutada (citas, usuarios, etc.)"],
        ["operacion", "VARCHAR", "10", "-", "Tipo de evento transaccional: INSERT, UPDATE, DELETE"],
        ["datos_anteriores", "JSONB", "Variable", "-", "Snapshot inmutable del estado previo del registro"],
        ["datos_nuevos", "JSONB", "Variable", "-", "Snapshot inmutable del nuevo estado registrado"]
    ]
    t_shape30 = s30.shapes.add_table(len(rows30) + 1, 5, Inches(0.92), Inches(2.35), Inches(11.50), Inches(4.55))
    t30 = t_shape30.table
    t30.columns[0].width = Inches(2.60)
    t30.columns[1].width = Inches(1.60)
    t30.columns[2].width = Inches(1.20)
    t30.columns[3].width = Inches(1.10)
    t30.columns[4].width = Inches(5.00)
    for c_idx, h in enumerate(headers_dict):
        c = t30.cell(0, c_idx)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.fill.solid()
        c.fill.fore_color.rgb = C_WHITE
        p = c.text_frame.paragraphs[0]
        p.text = h
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = C_DARK
        p.alignment = PP_ALIGN.CENTER

    for r_idx, r_data in enumerate(rows30):
        bg = C_WHITE if r_idx % 2 == 0 else C_LIGHT_BG
        for c_idx, val in enumerate(r_data):
            c = t30.cell(r_idx + 1, c_idx)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            c.fill.fore_color.rgb = bg
            p = c.text_frame.paragraphs[0]
            p.text = val
            p.font.name = "Arial"
            p.font.size = Pt(10.5)
            p.font.color.rgb = C_DARK
            if c_idx == 0:
                p.font.bold = True
            elif c_idx in [1, 2, 3]:
                p.alignment = PP_ALIGN.CENTER
                if c_idx == 3 and "PK" in val:
                    p.font.bold = True
    set_speaker_notes(s30, "Diccionario de catálogo SLA y esquema JSONB inmutable para el cumplimiento del requerimiento de auditoría RNF-03.")

    # =========================================================================
    # SECCIÓN 4: DESARROLLO (SLIDE 32: SOFTWARE USADO)
    # =========================================================================
    s32 = prs.slides[31]
    set_slide_title(s32, "Software Usado")
    clean_slide_body(s32)

    sub32 = s32.shapes.add_textbox(Inches(0.92), Inches(1.80), Inches(11.50), Inches(0.45))
    p_sub32 = sub32.text_frame.paragraphs[0]
    p_sub32.text = "Stack Tecnológico y Herramientas Empleadas en el Desarrollo de ServiTech:"
    p_sub32.font.name = "Arial"
    p_sub32.font.size = Pt(15.5)
    p_sub32.font.bold = True
    p_sub32.font.color.rgb = C_DARK

    headers32 = ["Capa de Arquitectura", "Tecnología / Software", "Versión", "Rol y Función en el Sistema"]
    rows32 = [
        ["Backend & Lógica", "Python", "3.12 LTS", "Lenguaje principal de desarrollo backend y procesamiento"],
        ["Framework Web", "Django", "5.1+", "Arquitectura MVT, ORM, motor de autenticación y sesiones"],
        ["Base de Datos", "PostgreSQL", "16.x", "Motor relacional ACID con triggers PL/pgSQL y soporte JSONB"],
        ["Frontend & UI", "HTML5 / CSS3", "Estándar W3C", "Estructura semántica, diseño accesible y responsivo"],
        ["Framework CSS", "TailwindCSS / Bootstrap", "3.4 / 5.3", "Estilos corporativos modernos, modo oscuro y microinteracciones"],
        ["Interactividad Web", "JavaScript ES6+", "Nativo", "Validaciones de formularios y actualización asíncrona de turnos"],
        ["Control de Versiones", "Git & GitHub", "2.4x", "Repositorio remoto, ramas de desarrollo y control de cambios"],
        ["Diseño & Prototipado", "Figma", "Cloud", "Diseño de interfaces, wireframes y experiencia de usuario (UX/UI)"]
    ]
    t_shape32 = s32.shapes.add_table(len(rows32) + 1, 4, Inches(0.92), Inches(2.35), Inches(11.50), Inches(4.55))
    t32 = t_shape32.table
    t32.columns[0].width = Inches(2.50)
    t32.columns[1].width = Inches(2.60)
    t32.columns[2].width = Inches(1.50)
    t32.columns[3].width = Inches(4.90)
    for c_idx, h in enumerate(headers32):
        c = t32.cell(0, c_idx)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.fill.solid()
        c.fill.fore_color.rgb = C_WHITE
        p = c.text_frame.paragraphs[0]
        p.text = h
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = C_DARK
        p.alignment = PP_ALIGN.CENTER

    for r_idx, r_data in enumerate(rows32):
        bg = C_WHITE if r_idx % 2 == 0 else C_LIGHT_BG
        for c_idx, val in enumerate(r_data):
            c = t32.cell(r_idx + 1, c_idx)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            c.fill.fore_color.rgb = bg
            p = c.text_frame.paragraphs[0]
            p.text = val
            p.font.name = "Arial"
            p.font.size = Pt(11)
            p.font.color.rgb = C_DARK
            if c_idx == 0:
                p.font.bold = True
            elif c_idx in [1, 2]:
                p.alignment = PP_ALIGN.CENTER
    set_speaker_notes(s32, "El stack combina la robustez y seguridad de Django y PostgreSQL con una interfaz moderna en TailwindCSS.")

    # =========================================================================
    # SECCIÓN 5: IMPLEMENTACIÓN (SLIDE 34: CASOS DE PRUEBA)
    # =========================================================================
    s34 = prs.slides[33]
    set_slide_title(s34, "Especificación de Caso de Prueba")
    clean_slide_body(s34)

    sub34 = s34.shapes.add_textbox(Inches(0.92), Inches(1.80), Inches(11.50), Inches(0.45))
    p_sub34 = sub34.text_frame.paragraphs[0]
    p_sub34.text = "Matriz de Ejecución y Validación de Casos de Prueba del Sistema:"
    p_sub34.font.name = "Arial"
    p_sub34.font.size = Pt(15.5)
    p_sub34.font.bold = True
    p_sub34.font.color.rgb = C_DARK

    headers34 = ["ID Caso", "Módulo / Función", "Entrada / Acción", "Resultado Esperado", "Resultado"]
    rows34 = [
        ("CP-01", "Motor de Agendamiento", "Cliente selecciona dispositivo, servicio y franja horaria disponible.", "Cita registrada en BD, franja horaria bloqueada y ticket generado.", "Aprobado", C_GREEN_PASS),
        ("CP-02", "Control de Colisiones", "Segundo cliente intenta reservar la misma franja ya ocupada.", "Sistema rechaza la reserva y muestra alerta de indisponibilidad.", "Aprobado", C_GREEN_PASS),
        ("CP-03", "Aviso 'Llegaré Tarde'", "Cliente pulsa botón de aviso (+15 min) desde recordatorio.", "Estado cambia a 'Retrasado con Aviso' y amplía tolerancia en box.", "Aprobado", C_GREEN_PASS),
        ("CP-04", "Pausa de Turno Técnico", "Técnico intenta pausar turno con orden activa en ejecución.", "Sistema bloquea la pausa exigiendo finalizar o liberar la orden.", "Aprobado", C_GREEN_PASS),
        ("CP-05", "Auditoría Automática", "Modificación o eliminación directa de registro en tabla citas.", "Trigger PL/pgSQL genera snapshot inmutable en auditoria_log (< 5 ms).", "Aprobado", C_GREEN_PASS),
        ("CP-06", "Consulta de Ticket", "Cliente ingresa código de orden en buscador de seguimiento público.", "Despliegue inmediato del estado actual del equipo sin requerir login.", "Aprobado", C_GREEN_PASS)
    ]
    t_shape34 = s34.shapes.add_table(len(rows34) + 1, 5, Inches(0.92), Inches(2.35), Inches(11.50), Inches(4.55))
    t34 = t_shape34.table
    t34.columns[0].width = Inches(1.20)
    t34.columns[1].width = Inches(2.60)
    t34.columns[2].width = Inches(3.20)
    t34.columns[3].width = Inches(3.30)
    t34.columns[4].width = Inches(1.20)
    for c_idx, h in enumerate(headers34):
        c = t34.cell(0, c_idx)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.fill.solid()
        c.fill.fore_color.rgb = C_WHITE
        p = c.text_frame.paragraphs[0]
        p.text = h
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = C_DARK
        p.alignment = PP_ALIGN.CENTER

    for r_idx, (cp_id, mod, accion, esperado, estado, col_res) in enumerate(rows34):
        bg = C_WHITE if r_idx % 2 == 0 else C_LIGHT_BG
        row_vals = [cp_id, mod, accion, esperado, estado]
        for c_idx, val in enumerate(row_vals):
            c = t34.cell(r_idx + 1, c_idx)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            c.fill.fore_color.rgb = bg
            p = c.text_frame.paragraphs[0]
            p.text = val
            p.font.name = "Arial"
            p.font.size = Pt(10.5)
            p.font.color.rgb = C_DARK
            if c_idx == 0:
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
            elif c_idx == 4:
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = col_res
    set_speaker_notes(s34, "Matriz de casos de prueba ejecutados satisfactoriamente, validando integridad, seguridad y reglas de negocio.")

    prs.save(pptx_path)
    print("Presentación completa de ServiTech actualizada exitosamente en Servitech-app.pptx!")

if __name__ == "__main__":
    update_entire_presentation()
