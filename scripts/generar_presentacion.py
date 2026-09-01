import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    C_GREEN = RGBColor(57, 169, 0)       # SENA Green (#39A900)
    C_NAVY = RGBColor(0, 50, 77)         # SENA Navy (#00324D)
    C_DARK = RGBColor(30, 41, 59)        # Slate Dark (#1E293B)
    C_GRAY = RGBColor(100, 116, 139)     # Slate Gray (#64748B)
    C_LIGHT = RGBColor(248, 250, 252)    # Slate Light (#F8FAFC)
    C_WHITE = RGBColor(255, 255, 255)
    C_CARD_BG = RGBColor(255, 255, 255)
    C_CARD_BORDER = RGBColor(226, 232, 240)
    C_BLUE = RGBColor(37, 99, 235)

    blank_layout = prs.slide_layouts[6]

    def set_slide_background(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, category_text=""):
        # Header banner
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        if category_text:
            p0 = tf.paragraphs[0]
            p0.text = category_text.upper()
            p0.font.name = "Arial"
            p0.font.size = Pt(11)
            p0.font.bold = True
            p0.font.color.rgb = C_GREEN
            p0.space_after = Pt(2)
            p1 = tf.add_paragraph()
        else:
            p1 = tf.paragraphs[0]

        p1.text = title_text
        p1.font.name = "Arial"
        p1.font.size = Pt(22)
        p1.font.bold = True
        p1.font.color.rgb = C_NAVY

        # Accent line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35), Inches(2.5), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = C_GREEN
        line.line.color.rgb = C_GREEN

        # Footer badge
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4))
        tf_f = footer_box.text_frame
        tf_f.margin_left = tf_f.margin_top = tf_f.margin_right = tf_f.margin_bottom = 0
        p_f = tf_f.paragraphs[0]
        p_f.text = "Servitech-app • SENA ADSO / Programación de Software • Septiembre 2026"
        p_f.font.name = "Arial"
        p_f.font.size = Pt(9)
        p_f.font.color.rgb = C_GRAY

    def add_speaker_notes(slide, text):
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = text

    def add_card(slide, left, top, width, height, title, items, badge=""):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = C_CARD_BG
        shape.line.color.rgb = C_CARD_BORDER
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.2)
        tf.margin_bottom = Inches(0.2)

        p = tf.paragraphs[0]
        if badge:
            p.text = f"{badge}  {title}"
        else:
            p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_NAVY
        p.space_after = Pt(10)

        for item in items:
            p_item = tf.add_paragraph()
            p_item.text = f"• {item}"
            p_item.font.name = "Arial"
            p_item.font.size = Pt(11)
            p_item.font.color.rgb = C_DARK
            p_item.space_after = Pt(6)

    def create_separator_slide(sec_num, sec_title, subtitle=""):
        slide = prs.slides.add_slide(blank_layout)
        set_slide_background(slide, C_NAVY)

        # Green accent block
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
        accent.fill.solid()
        accent.fill.fore_color.rgb = C_GREEN
        accent.line.fill.background()

        # Text container
        tb = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.5), Inches(3.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = f"SECCIÓN {sec_num}"
        p0.font.name = "Arial"
        p0.font.size = Pt(16)
        p0.font.bold = True
        p0.font.color.rgb = C_GREEN
        p0.space_after = Pt(10)

        p1 = tf.add_paragraph()
        p1.text = sec_title
        p1.font.name = "Arial"
        p1.font.size = Pt(36)
        p1.font.bold = True
        p1.font.color.rgb = C_WHITE
        p1.space_after = Pt(10)

        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.name = "Arial"
            p2.font.size = Pt(16)
            p2.font.color.rgb = RGBColor(203, 213, 225)

        return slide

    def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
        num_rows = len(rows) + 1
        num_cols = len(headers)
        table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
        table = table_shape.table

        if col_widths:
            for idx, w in enumerate(col_widths):
                table.columns[idx].width = w

        # Format Headers
        for col_idx, header_text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_NAVY
            p = cell.text_frame.paragraphs[0]
            p.text = header_text
            p.font.name = "Arial"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = C_WHITE
            p.alignment = PP_ALIGN.CENTER

        # Format Data Rows
        for row_idx, row_data in enumerate(rows):
            bg_color = C_WHITE if row_idx % 2 == 0 else C_LIGHT
            for col_idx, cell_value in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_color
                p = cell.text_frame.paragraphs[0]
                p.text = str(cell_value)
                p.font.name = "Arial"
                p.font.size = Pt(10)
                p.font.color.rgb = C_DARK

    # ==========================================
    # SLIDE 1: Portada
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, C_NAVY)

    # Accent shape
    accent1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.5), Inches(7.5))
    accent1.fill.solid()
    accent1.fill.fore_color.rgb = C_GREEN
    accent1.line.fill.background()

    tb1 = slide1.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.5), Inches(4.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "SISTEMA WEB DE AGENDAMIENTO Y GESTIÓN DE SERVICIOS DE REPARACIÓN TÉCNICA"
    p.font.name = "Arial"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.space_after = Pt(12)

    p2 = tf1.add_paragraph()
    p2.text = "Servitech-app • Proyecto Final de Sustentación"
    p2.font.name = "Arial"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = C_GREEN
    p2.space_after = Pt(20)

    p3 = tf1.add_paragraph()
    p3.text = "Servicio Nacional de Aprendizaje (SENA) | Clasificación: Pública reservada"
    p3.font.name = "Arial"
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(148, 163, 184)

    add_speaker_notes(slide1, "Buenos días estimados instructores y jurados. Presentamos Servitech-app, una solución de software web desarrollada para centralizar, optimizar y automatizar el agendamiento y la gestión de servicios de reparación técnica.")

    # ==========================================
    # SLIDE 2: Aprendices
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, C_LIGHT)
    add_header(slide2, "Equipo de Trabajo y Aprendices", "SENA - ADSO / Programación de Software")

    add_card(slide2, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Aprendices Investigadores / Desarrolladores", [
        "Juan Bayona — Análisis de Requerimientos & QA",
        "Cristian Contreras — Arquitectura & Base de Datos",
        "Emmanuel Flores — Diseño UI/UX & Frontend",
        "Breyner Peña — Backend Django & Auditoría PostgreSQL"
    ], "👥")

    add_card(slide2, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Información del Proyecto", [
        "Programa: Análisis y Desarrollo de Software (ADSO)",
        "Centro de Formación: SENA Regional",
        "Convocatoria / Sustentación: Septiembre 2026",
        "Stack: Python (Django), PostgreSQL (PL/pgSQL), Tailwind CSS, JS",
        "Control de Tareas: Jira Software & Git / GitHub"
    ], "📋")

    add_speaker_notes(slide2, "Nuestro equipo de trabajo está integrado por Juan Bayona, Cristian Contreras, Emmanuel Flores y Breyner Peña. A lo largo de la presentación demostraremos la rigurosidad técnica aplicada en cada fase del ciclo de vida del software.")

    # ==========================================
    # SLIDE 3: Separador I
    # ==========================================
    slide3 = create_separator_slide("I", "FORMULACIÓN DEL PROYECTO", "Problema, Justificación, Objetivos, Alcance, Riesgos y Cronograma")
    add_speaker_notes(slide3, "Iniciamos con la Formulación del Proyecto, donde establecemos la problemática de negocio y el alcance de la solución propuesta.")

    # ==========================================
    # SLIDE 4: Planteamiento del Problema
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, C_LIGHT)
    add_header(slide4, "Planteamiento del Problema", "Sección I: Formulación")

    add_card(slide4, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), "Gestión Manual y Desorden", [
        "Uso de planillas en papel y cuadernos de turnos.",
        "Pérdida de historial de diagnósticos previos.",
        "Falta de trazabilidad sobre componentes cambiados."
    ], "⚠️")

    add_card(slide4, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8), "Incertidumbre del Cliente", [
        "Poca visibilidad del estado real del dispositivo.",
        "Tiempos de espera excesivos en taller presencial.",
        "Dificultades en la validez y reclamo de garantías."
    ], "⏳")

    add_card(slide4, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.8), "Fugas e Inasistencias", [
        "Pérdida y descontrol en repuestos utilizados.",
        "Ausentismo de clientes (No-Show) que congela boxes.",
        "Falta de datos consolidados para la gerencia."
    ], "📉")

    add_speaker_notes(slide4, "La gestión manual en talleres de reparación tecnológica causa desorden, pérdidas de repuestos y descontento en los clientes. Servitech-app nace para erradicar estas ineficiencias.")

    # ==========================================
    # SLIDE 5: Justificación
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, C_LIGHT)
    add_header(slide5, "Justificación del Proyecto", "Sección I: Formulación")

    add_card(slide5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Impacto Operativo y Económico", [
        "Optimización de tiempos: Reducción del 60% en tiempos de espera presencial.",
        "Control de repuestos: Trazabilidad exacta de cada componente instalado en una orden.",
        "Aprovechamiento técnico: Asignación por especialidad (Nivel 1, 2 y 3)."
    ], "🚀")

    add_card(slide5, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Valor Agregado para el Cliente", [
        "Autonomía total: Agendamiento online 24/7 en solo 4 pasos.",
        "Transparencia: Consulta del estado de reparación en tiempo real.",
        "Gestión de contingencias: Aviso de retraso en 1 clic para no perder el cupo."
    ], "⭐")

    add_speaker_notes(slide5, "La plataforma optimiza el flujo de trabajo, agiliza la comunicación cliente-taller y garantiza que cada orden de trabajo tenga soporte en inventario y auditoría.")

    # ==========================================
    # SLIDE 6: Objetivo General
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, C_LIGHT)
    add_header(slide6, "Objetivo General", "Sección I: Formulación")

    shape_obj = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(2.2), Inches(10.9), Inches(3.8))
    shape_obj.fill.solid()
    shape_obj.fill.fore_color.rgb = C_NAVY
    shape_obj.line.color.rgb = C_GREEN
    shape_obj.line.width = Pt(2.5)

    tf = shape_obj.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_right = Inches(0.6)
    tf.margin_top = Inches(0.6)

    p = tf.paragraphs[0]
    p.text = "OBJETIVO GENERAL DEL PROYECTO"
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = C_GREEN
    p.space_after = Pt(14)

    p1 = tf.add_paragraph()
    p1.text = "Desarrollar un sistema web centralizado para la gestión integral de órdenes de servicio, agendamiento de citas y control de inventario en talleres de soporte técnico, optimizando el flujo de trabajo, la trazabilidad de los dispositivos y la toma de decisiones gerenciales."
    p1.font.name = "Arial"
    p1.font.size = Pt(18)
    p1.font.color.rgb = C_WHITE

    add_speaker_notes(slide6, "Nuestro objetivo general abarca desde la captura ágil del dispositivo hasta el control de inventario y la analítica ejecutiva.")

    # ==========================================
    # SLIDE 7: Objetivos Específicos
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, C_LIGHT)
    add_header(slide7, "Objetivos Específicos", "Sección I: Formulación")

    add_card(slide7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.3), "1. Módulo de Recepción y Dispositivos", [
        "Registrar equipos, seriales/IMEI, fallas y datos del cliente con validaciones ágiles."
    ], "📱")

    add_card(slide7, Inches(6.8), Inches(1.8), Inches(5.6), Inches(2.3), "2. Automatización de Estados en Tiempo Real", [
        "Actualizar el avance técnico y notificar imprevistos o retrasos de forma bidireccional."
    ], "🔄")

    add_card(slide7, Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.3), "3. Control de Inventario y Repuestos", [
        "Vincular la salida de componentes a cada orden de servicio, controlando stock crítico."
    ], "📦")

    add_card(slide7, Inches(6.8), Inches(4.3), Inches(5.6), Inches(2.3), "4. Auditoría y Reportes Gerenciales", [
        "Generar reportes en Excel/PDF y registrar auditoría nativa con triggers PL/pgSQL."
    ], "📊")

    add_speaker_notes(slide7, "Definimos cuatro objetivos específicos directamente mapeados a los módulos de software construidos en la plataforma.")

    # ==========================================
    # SLIDE 8: Alcance
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, C_LIGHT)
    add_header(slide8, "Alcance del Sistema", "Sección I: Formulación")

    add_card(slide8, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), "Módulo Cliente", [
        "Registro y gestión de perfil.",
        "Asistente de citas en 4 pasos.",
        "Selección de técnico y horario.",
        "Aviso de retraso y reagendamiento."
    ], "🧑‍💻")

    add_card(slide8, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8), "Módulo Técnico", [
        "Dashboard con agenda diaria.",
        "Switch de estados de reparación.",
        "Pausa de turno y descansos.",
        "Descarga de repuestos utilizados."
    ], "🔧")

    add_card(slide8, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.8), "Módulo Administrativo", [
        "Gestión de usuarios y técnicos.",
        "Catálogo de servicios y SLA.",
        "Kardex e inventario de piezas.",
        "Exportación de reportes PDF/Excel."
    ], "🛡️")

    add_speaker_notes(slide8, "El alcance cubre los procesos internos del taller y la interacción de agendamiento autónomo para clientes finales.")

    # ==========================================
    # SLIDE 9: Riesgos y Mitigación
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9, C_LIGHT)
    add_header(slide9, "Matriz de Riesgos y Mitigación", "Sección I: Formulación")

    headers9 = ["Riesgo Identificado", "Impacto", "Estrategia de Mitigación en el Software"]
    rows9 = [
        ["Caídas de Red o Servidor", "Alto", "Persistencia de sesiones, base de datos transaccional ACID en PostgreSQL."],
        ["Resistencia al Cambio", "Medio", "Interfaces limpias con Tailwind CSS, botones grandes y flujos intuitivos."],
        ["Omisión Operativa de Técnicos", "Medio", "Alertas en dashboard, requerimiento de diagnóstico antes de finalizar orden."],
        ["Inasistencia de Clientes (No-Show)", "Medio", "Tolerancia de 15 min y botón de liberación rápida para clientes presenciales."]
    ]
    add_table(slide9, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.6), headers9, rows9, [Inches(3.2), Inches(1.5), Inches(7.0)])
    add_speaker_notes(slide9, "Se evaluaron los riesgos técnicos y operativos, implementando mecanismos preventivos en la lógica de negocio.")

    # ==========================================
    # SLIDE 10: Restricciones
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10, C_LIGHT)
    add_header(slide10, "Restricciones del Sistema", "Sección I: Formulación")

    add_card(slide10, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Restricciones Tecnológicas", [
        "Backend obligatorio: Python 3.12 / Django 6.0.",
        "Base de datos obligatoria: PostgreSQL con triggers PL/pgSQL.",
        "Frontend: HTML5, JavaScript moderno y Tailwind CSS.",
        "Control de tareas: Metodología ágil gestionada en Jira Software."
    ], "⚙️")

    add_card(slide10, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Restricciones de Hardware y Entorno", [
        "Equipos de cómputo en taller: Mínimo 4 GB de RAM.",
        "Compatibilidad: Navegadores web modernos (Chrome, Edge, Firefox).",
        "Conectividad: Red local o enlace a servidor en la nube.",
        "Normativa: Estándares de calidad de software y directrices SENA."
    ], "🖥️")

    add_speaker_notes(slide10, "El desarrollo se rigió por las restricciones de arquitectura, hardware mínimo y herramientas de gestión aprobadas.")

    # ==========================================
    # SLIDE 11: Cronograma
    # ==========================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11, C_LIGHT)
    add_header(slide11, "Cronograma de Ejecución", "Sección I: Formulación")

    headers11 = ["Fase", "Actividades Principales", "Entregables Clave", "Estado"]
    rows11 = [
        ["Fase 1: Análisis", "Levantamiento, historias de usuario, priorización MoSCoW", "Documento RF/RNF, Matriz de Trazabilidad", "Completada"],
        ["Fase 2: Diseño", "Modelado UML, diseño de base de datos DER, mockups", "Diagramas de Casos de Uso, Diccionario de Datos", "Completada"],
        ["Fase 3: Desarrollo", "Codificación Django MVT, triggers PostgreSQL, Tailwind", "Módulos de Agendamiento, Técnico y Admin", "Completada"],
        ["Fase 4: Pruebas", "Casos de prueba QA, validación de auditoría y reportes", "Matriz de Pruebas, Presentación de Sustentación", "Completada"]
    ]
    add_table(slide11, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.6), headers11, rows11, [Inches(2.2), Inches(4.5), Inches(3.5), Inches(1.5)])
    add_speaker_notes(slide11, "El proyecto cumplió satisfactoriamente con cada una de las 4 fases del ciclo de vida del software.")

    # ==========================================
    # SLIDE 12: Separador II
    # ==========================================
    slide12 = create_separator_slide("II", "ANÁLISIS DEL SISTEMA", "Técnicas de Captura, Stakeholders, Historias de Usuario, Priorización, RF y RNF")
    add_speaker_notes(slide12, "Pasamos a la Sección II de Análisis, donde exponemos los requerimientos funcionales, no funcionales y la priorización técnica.")

    # ==========================================
    # SLIDE 13: Herramientas Captura Requisitos
    # ==========================================
    slide13 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide13, C_LIGHT)
    add_header(slide13, "Técnicas para la Captura de Requisitos", "Sección II: Análisis")

    add_card(slide13, Inches(0.8), Inches(1.8), Inches(2.7), Inches(4.8), "Entrevistas", [
        "Aplicadas a administradores y técnicos.",
        "Identificación de tiempos muertos y reglas de asignación.",
        "Definición de niveles de especialidad (N1, N2, N3)."
    ], "🗣️")

    add_card(slide13, Inches(3.8), Inches(1.8), Inches(2.7), Inches(4.8), "Encuestas", [
        "Muestra de 45 clientes habituales de talleres.",
        "88% insatisfecho con la falta de seguimiento.",
        "Preferencia por agendamiento autónomo 24/7."
    ], "📋")

    add_card(slide13, Inches(6.8), Inches(1.8), Inches(2.7), Inches(4.8), "Observación Directa", [
        "Acompañamiento en taller presencial.",
        "Mapeo del ciclo de vida del dispositivo.",
        "Detección de pérdidas en consumo de repuestos."
    ], "👁️")

    add_card(slide13, Inches(9.8), Inches(1.8), Inches(2.7), Inches(4.8), "Análisis de Artefactos", [
        "Revisión de boletas físicas de recepción.",
        "Formatos de garantía y planillas de Excel.",
        "Normalización hacia entidades relacionales."
    ], "📑")

    add_speaker_notes(slide13, "La captura de requisitos se fundamentó en fuentes primarias: entrevistas a técnicos, encuestas a clientes y análisis de planillas físicas.")

    # ==========================================
    # SLIDE 14: Matriz de Stakeholders
    # ==========================================
    slide14 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide14, C_LIGHT)
    add_header(slide14, "Matriz de Stakeholders (Partes Interesadas)", "Sección II: Análisis")

    headers14 = ["Stakeholder", "Tipo", "Interés", "Poder", "Expectativa Principal", "Estrategia"]
    rows14 = [
        ["Administrador", "Interno", "Alto", "Alto", "Control global, auditoría, reportes ejecutivos", "Gestionar de cerca"],
        ["Técnico de Taller", "Interno", "Alto", "Medio", "Agenda clara, registro de piezas, switch estados", "Mantener satisfecho"],
        ["Cliente Final", "Externo", "Alto", "Bajo", "Agendamiento ágil, consulta y aviso de retraso", "Mantener informado"],
        ["Recepcionista", "Interno", "Medio", "Medio", "Registro rápido de walk-ins y control No-Show", "Mantener informado"],
        ["Comité SENA", "Externo", "Alto", "Alto", "Cumplimiento de estándares y calidad de software", "Gestionar de cerca"]
    ]
    add_table(slide14, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.6), headers14, rows14, [Inches(2.0), Inches(1.2), Inches(1.2), Inches(1.2), Inches(4.1), Inches(2.0)])
    add_speaker_notes(slide14, "La matriz de stakeholders clasifica el nivel de interés y poder de cada rol para orientar el diseño del sistema a sus necesidades.")

    # ==========================================
    # SLIDE 15: Historias de Usuario
    # ==========================================
    slide15 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide15, C_LIGHT)
    add_header(slide15, "Historias de Usuario (User Stories)", "Sección II: Análisis")

    add_card(slide15, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.3), "HU-01: Selección de Equipo y Servicio", [
        "Como: Cliente final.",
        "Quiero: Elegir mi dispositivo y tipo de reparación.",
        "Para: Reservar la duración adecuada según el catálogo."
    ], "🧑‍💻")

    add_card(slide15, Inches(6.8), Inches(1.8), Inches(5.6), Inches(2.3), "HU-03: Aviso de Retraso", [
        "Como: Cliente agendado.",
        "Quiero: Notificar un retraso (+10/+15 min) en 1 clic.",
        "Para: Conservar mi cupo sin cancelación automática."
    ], "⏱️")

    add_card(slide15, Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.3), "HU-06: Panel Operativo del Técnico", [
        "Como: Técnico asignado.",
        "Quiero: Ver mis citas del día y cambiar estados.",
        "Para: Organizar mi mesa de trabajo y tiempos de atención."
    ], "🔧")

    add_card(slide15, Inches(6.8), Inches(4.3), Inches(5.6), Inches(2.3), "HU-10: Auditoría y Reportes PDF", [
        "Como: Administrador.",
        "Quiero: Exportar KPIs en PDF y auditar operaciones.",
        "Para: Tomar decisiones estratégicas con datos íntegros."
    ], "📊")

    add_speaker_notes(slide15, "Las Historias de Usuario se redactaron bajo la estructura ágil estándar con criterios de aceptación claros y medibles.")

    # ==========================================
    # SLIDE 16: Matriz de Trazabilidad HU vs RF
    # ==========================================
    slide16 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide16, C_LIGHT)
    add_header(slide16, "Matriz de Trazabilidad (HU vs. Requerimientos)", "Sección II: Análisis")

    headers16 = ["Historia (HU)", "Título de la Historia", "Requerimiento Funcional (RF)", "Módulo", "Estado"]
    rows16 = [
        ["HU-01, HU-02", "Selección de Dispositivo, Servicio y Horario", "RF04 (Dispositivos), RF05 (Catálogo), RF06 (Citas)", "Agendamiento", "Implementado"],
        ["HU-03, HU-04", "Aviso de Retraso y Reagendamiento", "RF06 (Horarios), RF07 (Contingencias)", "Turnos / Citas", "Implementado"],
        ["HU-05, HU-06", "Ajuste de Tiempo y Panel del Técnico", "RF08 (Panel Técnico), RF09 (Historial Cita)", "Operación Técnica", "Implementado"],
        ["HU-07, HU-08", "Control No-Show y Parametrización SLA", "RF03 (Usuarios), RF05 (Servicios), RF10 (Pausas)", "Administración", "Implementado"],
        ["HU-09, HU-10", "Inventario por Orden y Auditoría / Reportes", "RF02 (Auditoría), RF11/12 (Inventario), RF14/15 (Reportes)", "Analítica / BD", "Implementado"]
    ]
    add_table(slide16, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.6), headers16, rows16, [Inches(1.8), Inches(3.7), Inches(3.2), Inches(1.8), Inches(1.2)])
    add_speaker_notes(slide16, "Garantizamos trazabilidad total: cada historia de usuario está vinculada a requerimientos funcionales y a código en producción.")

    # ==========================================
    # SLIDE 17: Priorización MoSCoW
    # ==========================================
    slide17 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide17, C_LIGHT)
    add_header(slide17, "Técnica de Priorización MoSCoW & Conclusión", "Sección II: Análisis")

    add_card(slide17, Inches(0.8), Inches(1.8), Inches(2.7), Inches(4.8), "Must Have (Obligatorio)", [
        "RF01: Autenticación y roles.",
        "RF06: Agendamiento sin solapamiento.",
        "RF08/09: Panel de técnico y estados.",
        "RF11/12: Repuestos e inventario.",
        "RF02/15: Auditoría y PDF."
    ], "🟢")

    add_card(slide17, Inches(3.8), Inches(1.8), Inches(2.7), Inches(4.8), "Should Have (Deseable)", [
        "RF07: Aviso de retraso de clientes.",
        "RF10: Pausa manual de turno.",
        "RF07: Liberación rápida por No-Show."
    ], "🟡")

    add_card(slide17, Inches(6.8), Inches(1.8), Inches(2.7), Inches(4.8), "Could Have (Opcional)", [
        "RF16: Módulo de tickets de soporte interno.",
        "Filtros avanzados de búsqueda en historial."
    ], "🔵")

    add_card(slide17, Inches(9.8), Inches(1.8), Inches(2.7), Inches(4.8), "Won't Have (Post-MVP)", [
        "RF17: Pasarela de pagos con tarjeta online.",
        "Integración con WhatsApp API oficial."
    ], "⚪")

    add_speaker_notes(slide17, "La priorización MoSCoW permitió enfocar los esfuerzos en un MVP completamente funcional y libre de dependencias no esenciales.")

    # ==========================================
    # SLIDE 18: Requerimientos Funcionales
    # ==========================================
    slide18 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide18, C_LIGHT)
    add_header(slide18, "Catálogo de Requerimientos Funcionales (RF)", "Sección II: Análisis")

    headers18 = ["Código", "Requerimiento Funcional", "Descripción Técnica", "Prioridad"]
    rows18 = [
        ["RF01", "Autenticación por Roles", "Acceso por correo y contraseña; redirección a Admin, Técnico o Cliente.", "Must Have"],
        ["RF04", "Registro de Dispositivos", "Captura de IMEI/Serial, marca, modelo y estado físico.", "Must Have"],
        ["RF06", "Motor de Agendamiento", "Reserva de bloques de tiempo sin colisiones de agenda.", "Must Have"],
        ["RF09", "Historial de Estados", "Registro inmutable del cambio de fases en cada orden de trabajo.", "Must Have"],
        ["RF12", "Kardex de Inventario", "Descuento automático de existencias al usar repuestos en citas.", "Must Have"],
        ["RF15", "Exportación PDF/Excel", "Generación server-side de informes analíticos con reportlab y openpyxl.", "Must Have"]
    ]
    add_table(slide18, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.6), headers18, rows18, [Inches(1.2), Inches(2.8), Inches(6.2), Inches(1.5)])
    add_speaker_notes(slide18, "Los Requerimientos Funcionales especifican las operaciones y transformaciones de datos que realiza el sistema.")

    # ==========================================
    # SLIDE 19: Requerimientos No Funcionales
    # ==========================================
    slide19 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide19, C_LIGHT)
    add_header(slide19, "Requerimientos No Funcionales (ISO/IEC 25010)", "Sección II: Análisis")

    headers19 = ["Código", "Dimensión ISO 25010", "Requerimiento No Funcional", "Criterio de Aceptación / Métrica"]
    rows19 = [
        ["RNF01", "Rendimiento y Eficiencia", "Tiempo de respuesta en consultas de agenda", "Tiempo de respuesta menor a 1.5 segundos."],
        ["RNF02", "Seguridad y Privacidad", "Protección de credenciales y prevención de ataques", "Hash PBKDF2-SHA256, tokens CSRF y ORM seguro."],
        ["RNF03", "Trazabilidad e Integridad", "Auditoría transaccional en base de datos", "Triggers nativos en PL/pgSQL hacia auditoria_log."],
        ["RNF04", "Usabilidad y Accesibilidad", "Diseño responsivo y amigable Mobile-First", "Tailwind CSS con cumplimiento de contraste WCAG AA."],
        ["RNF05", "Confiabilidad (ACID)", "Consistencia en transacciones concurrentes", "Integridad referencial y bloqueo atómico de turnos."]
    ]
    add_table(slide19, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.6), headers19, rows19, [Inches(1.2), Inches(2.6), Inches(4.2), Inches(3.7)])
    add_speaker_notes(slide19, "Los RNF siguen la norma ISO/IEC 25010 para asegurar altos niveles de calidad, seguridad y velocidad de respuesta.")

    # ==========================================
    # SLIDE 20: Separador III
    # ==========================================
    slide20 = create_separator_slide("III", "DISEÑO DEL SISTEMA", "Scrum, Casos de Uso UML, Mockups, Arquitectura, Base de Datos y Diccionario")
    add_speaker_notes(slide20, "Entramos a la Sección III de Diseño, que detalla la arquitectura de software, modelado de base de datos y diseño de interfaces.")

    # ==========================================
    # SLIDE 21: Metodología Scrum y Jira
    # ==========================================
    slide21 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide21, C_LIGHT)
    add_header(slide21, "Metodología de Desarrollo Ágil: Scrum", "Sección III: Diseño")

    add_card(slide21, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Marco de Trabajo Scrum", [
        "Sprints de 1 a 2 semanas con entregables funcionales.",
        "Daily Standups para control de impedimentos.",
        "Sprint Planning y Sprint Review con demostraciones reales.",
        "Retrospectivas para mejora continua del código."
    ], "🏃")

    add_card(slide21, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Control de Tareas en Jira Software", [
        "Épicas estructuradas por módulos del sistema.",
        "Historias de usuario con estimación en Story Points.",
        "Tablero Kanban/Scrum para seguimiento del flujo.",
        "Vinculación de commits y pull requests en GitHub."
    ], "📊")

    add_speaker_notes(slide21, "Scrum y Jira nos permitieron organizar el desarrollo por sprints, manteniendo visibilidad y control sobre cada tarea.")

    # ==========================================
    # SLIDE 22: Diagrama de Casos de Uso UML
    # ==========================================
    slide22 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide22, C_LIGHT)
    add_header(slide22, "Diagrama General de Casos de Uso (UML)", "Sección III: Diseño")

    add_card(slide22, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), "Actor: Cliente", [
        "CU02: Agendar Cita Técnica.",
        "CU03: Notificar Retraso (+10/+15m).",
        "CU_Reg: Registrar Dispositivo.",
        "CU_Reag: Reagendar / Cancelar Cita."
    ], "🧑‍💻")

    add_card(slide22, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8), "Actor: Técnico", [
        "CU04: Gestionar Turno y Pausa.",
        "CU_At: Aceptar / Finalizar Cita.",
        "CU_Rep: Registrar Repuestos Usados.",
        "CU_Hist: Consultar Historial Cliente."
    ], "🔧")

    add_card(slide22, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.8), "Actor: Administrador", [
        "CU01: Visualizar Dashboard KPIs.",
        "CU05: Gestionar Catálogo y SLA.",
        "CU06: Exportar Reporte PDF/Excel.",
        "CU_Aud: Auditar Log de Operaciones."
    ], "🛡️")

    add_speaker_notes(slide22, "El modelado UML organiza las interacciones de los tres actores con los casos de uso del sistema.")

    # ==========================================
    # SLIDE 23: Tablas de Casos de Uso
    # ==========================================
    slide23 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide23, C_LIGHT)
    add_header(slide23, "Especificación de Casos de Uso (Tablas CU)", "Sección III: Diseño")

    headers23 = ["Caso de Uso", "Actor Principal", "Precondición", "Flujo Principal de Éxito"]
    rows23 = [
        ["CU02 - Agendar Cita", "Cliente", "Servicio activo en catálogo", "1. Elige equipo -> 2. Elige servicio -> 3. Selecciona horario/técnico -> 4. Confirma reserva."],
        ["CU03 - Demoras de Cita", "Técnico / Cliente", "Cita registrada en sistema", "1. Ingresa a cita -> 2. Pulsa retraso o ajuste -> 3. Sistema actualiza estado en panel."],
        ["CU04 - Pausa Técnica", "Técnico", "Sesión iniciada en panel", "1. Pulsa 'Pausar Turno' -> 2. Sistema verifica citas activas -> 3. Pasa a estado No Disponible."],
        ["CU06 - Exportar Reporte", "Administrador", "Datos en el periodo elegido", "1. Filtra rango -> 2. Pulsa 'Exportar PDF' -> 3. reportlab compila y descarga documento."]
    ]
    add_table(slide23, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.6), headers23, rows23, [Inches(2.5), Inches(1.8), Inches(2.4), Inches(5.0)])
    add_speaker_notes(slide23, "Cada caso de uso está especificado formalmente con precondiciones, flujo principal y flujos alternativos de error.")

    # ==========================================
    # SLIDE 24: Mockups y UI/UX
    # ==========================================
    slide24 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide24, C_LIGHT)
    add_header(slide24, "Mockups y Diseño de Interfaces (UI/UX)", "Sección III: Diseño")

    add_card(slide24, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), "Asistente de Agendamiento", [
        "Diseño paso a paso (Stepper).",
        "Tarjetas interactivas de dispositivos.",
        "Selector visual de bloques de tiempo.",
        "Confirmación con resumen detallado."
    ], "📱")

    add_card(slide24, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8), "Dashboard del Técnico", [
        "Vista de tarjetas operativas del día.",
        "Indicador visual de estado de turno.",
        "Badges de alerta para clientes retrasados.",
        "Botones directos 'Aceptar' y 'Finalizar'."
    ], "💻")

    add_card(slide24, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.8), "Panel Administrativo", [
        "Gráficos analíticos con Chart.js.",
        "Tablas dinámicas de inventario y stock.",
        "Gestión modular de usuarios y técnicos.",
        "Exportación de reportes a un clic."
    ], "📊")

    add_speaker_notes(slide24, "Los mockups priorizan la simplicidad visual, permitiendo que las operaciones críticas se completen rápidamente.")

    # ==========================================
    # SLIDE 25: Mapa de Navegación y Accesibilidad
    # ==========================================
    slide25 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide25, C_LIGHT)
    add_header(slide25, "Mapa de Navegación, Usabilidad y Accesibilidad", "Sección III: Diseño")

    add_card(slide25, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Arquitectura de Navegación", [
        "Rutas protegidas por middleware de rol en Django.",
        "Rutas públicas: Inicio, Login, Registro, Turno público.",
        "Rutas Cliente: `/cliente/inicio/`, `/servicios/dispositivo/`.",
        "Rutas Técnico: `/tecnico/`, `/tecnico/agenda/`, `/historial/`.",
        "Rutas Admin: `/admin-panel/citas/`, `/reportes/`, `/inventario/`."
    ], "🗺️")

    add_card(slide25, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Estándares de Usabilidad y a11y", [
        "Contraste adecuado en textos y botones (WCAG 2.1 AA).",
        "Navegación asistida por teclado en formularios.",
        "Componentes modales y Toast con feedback sonoro/visual.",
        "Diseño totalmente adaptable a pantallas móviles y tablets."
    ], "♿")

    add_speaker_notes(slide25, "La navegación está blindada por roles de seguridad y cumple con criterios internacionales de accesibilidad web.")

    # ==========================================
    # SLIDE 26: Modelo Entidad-Relación (DER)
    # ==========================================
    slide26 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide26, C_LIGHT)
    add_header(slide26, "Modelo Entidad-Relación (Base de Datos)", "Sección III: Diseño")

    headers26 = ["Entidad", "Relaciones Principales", "Cardinalidad", "Función en el Sistema"]
    rows26 = [
        ["Usuario", "Cita, Dispositivo, PerfilTecnico, AuditoriaLog", "1 a N / 1 a 1", "Gestión de cuentas, clientes y técnicos."],
        ["Dispositivo", "Usuario (Cliente), Cita", "N a 1 / 1 a N", "Registro de hardware con IMEI/Serial."],
        ["Cita", "Usuario (Cliente/Técnico), Servicio, Dispositivo, CitaRepuesto", "N a 1 / 1 a N", "Entidad central de agendamiento y turno."],
        ["Repuesto", "Inventario, CitaRepuesto", "1 a N", "Catálogo de piezas e insumos técnicos."],
        ["AuditoriaLog", "Disparada por triggers en PostgreSQL", "N/A", "Bitácora inmutable de cambios transaccionales."]
    ]
    add_table(slide26, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.6), headers26, rows26, [Inches(1.8), Inches(3.8), Inches(1.8), Inches(4.3)])
    add_speaker_notes(slide26, "La base de datos relacional está en Tercera Forma Normal (3FN), garantizando integridad referencial y atomicidad.")

    # ==========================================
    # SLIDE 27: Diccionario de Datos
    # ==========================================
    slide27 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide27, C_LIGHT)
    add_header(slide27, "Diccionario de Datos del Sistema", "Sección III: Diseño")

    headers27 = ["Tabla", "Campo", "Tipo de Dato", "Nulidad / Default", "Descripción"]
    rows27 = [
        ["turnos_usuario", "id_usuario", "BIGINT PRIMARY KEY", "NOT NULL / Auto", "Identificador único de usuario."],
        ["turnos_usuario", "rol", "VARCHAR(30)", "DEFAULT 'CLIENTE'", "Rol (ADMINISTRADOR, TECNICO, CLIENTE)."],
        ["turnos_cita", "fecha / hora_inicio", "DATE / TIME", "NOT NULL", "Día y hora de la cita programada."],
        ["turnos_cita", "minutos_retraso", "INTEGER", "DEFAULT 0", "Tiempo de tolerancia reportado por cliente."],
        ["turnos_repuesto", "stock / precio", "INT / DECIMAL(10,2)", "DEFAULT 0", "Existencias físicas y precio unitario."],
        ["auditoria_log", "datos_antes / despues", "JSONB", "NULLABLE", "Snapshot diferencial capturado por trigger."]
    ]
    add_table(slide27, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.6), headers27, rows27, [Inches(2.0), Inches(2.2), Inches(2.3), Inches(2.0), Inches(3.2)])
    add_speaker_notes(slide27, "El diccionario de datos define formalmente los atributos, tipos de datos y restricciones de integridad del modelo.")

    # ==========================================
    # SLIDE 28: Arquitectura de Software
    # ==========================================
    slide28 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide28, C_LIGHT)
    add_header(slide28, "Arquitectura de Software (Patrón MVT)", "Sección III: Diseño")

    add_card(slide28, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), "Modelos (M)", [
        "Mapeo Objeto-Relacional (ORM Django).",
        "Modelos modulares: Citas, Usuarios, Dispositivos, Inventario.",
        "Reglas de validación y métodos del dominio."
    ], "🗄️")

    add_card(slide28, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8), "Vistas / Controladores (V)", [
        "Controladores de lógica de negocio.",
        "Protección por decoradores y middleware.",
        "Despacho de respuestas HTML y JSON API.",
        "Exportadores con reportlab y openpyxl."
    ], "🎮")

    add_card(slide28, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.8), "Templates (T) & PL/pgSQL", [
        "Plantillas dinámicas con Tailwind CSS.",
        "Renderizado server-side seguro contra XSS.",
        "Triggers nativos en PostgreSQL para auditoría automática de transacciones."
    ], "🎨")

    add_speaker_notes(slide28, "La arquitectura MVT de Django desacopla la presentación de la lógica de negocio y del motor de base de datos.")

    # ==========================================
    # SLIDE 29: Stack Tecnológico y Calidad
    # ==========================================
    slide29 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide29, C_LIGHT)
    add_header(slide29, "Stack Tecnológico y Herramientas de Calidad", "Sección III: Diseño")

    add_card(slide29, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Tecnologías de Implementación", [
        "Lenguaje & Backend: Python 3.12 / Django 6.0.",
        "Base de Datos: PostgreSQL 16 con PL/pgSQL.",
        "Frontend: HTML5, JavaScript ES6, Tailwind CSS.",
        "Librerías clave: ReportLab (PDF), OpenPyXL (Excel), Chart.js."
    ], "🛠️")

    add_card(slide29, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Aseguramiento de Calidad", [
        "Control de versiones: Git & GitHub.",
        "Gestión ágil: Jira Software.",
        "Pruebas de endpoints: Postman.",
        "Estandarización y linting: ESLint / Flake8."
    ], "✨")

    add_speaker_notes(slide29, "Utilizamos un ecosistema moderno validado con herramientas de linting, pruebas en Postman y control de versiones en GitHub.")

    # ==========================================
    # SLIDE 30: Separador IV
    # ==========================================
    slide30 = create_separator_slide("IV", "DESARROLLO E IMPLEMENTACIÓN", "Módulos Codificados, Pruebas QA, Conclusiones y Demostración")
    add_speaker_notes(slide30, "Pasamos a la sección final de Desarrollo e Implementación, donde demostramos el software funcional y los resultados de las pruebas.")

    # ==========================================
    # SLIDE 31: Módulos Desarrollados
    # ==========================================
    slide31 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide31, C_LIGHT)
    add_header(slide31, "Módulos Desarrollados y Codificación", "Sección IV: Desarrollo")

    add_card(slide31, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.3), "Módulo de Agendamiento", [
        "Flujo en 4 pasos con cálculo dinámico de franjas horarias y prevención de solapamientos."
    ], "📅")

    add_card(slide31, Inches(6.8), Inches(1.8), Inches(5.6), Inches(2.3), "Módulo Técnico y Pausas", [
        "Panel interactivo de citas del día con switch de estados y botón de pausa manual de turno."
    ], "🔧")

    add_card(slide31, Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.3), "Módulo de Inventario & Repuestos", [
        "Kardex de movimientos y descarga automática de piezas utilizadas en reparaciones."
    ], "📦")

    add_card(slide31, Inches(6.8), Inches(4.3), Inches(5.6), Inches(2.3), "Auditoría & Reportes Ejecutivos", [
        "Triggers automáticos en PostgreSQL y exportador server-side de reportes analíticos a PDF."
    ], "📊")

    add_speaker_notes(slide31, "Todos los módulos han sido desarrollados, conectados al modelo de datos y probados en entorno real.")

    # ==========================================
    # SLIDE 32: Casos de Prueba QA
    # ==========================================
    slide32 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide32, C_LIGHT)
    add_header(slide32, "Casos de Prueba y Aseguramiento de Calidad (QA)", "Sección IV: Desarrollo")

    headers32 = ["ID", "Escenario de Prueba", "Procedimiento", "Resultado Esperado", "Estado"]
    rows32 = [
        ["CP01", "Autenticación por Roles", "Ingreso con credenciales de Admin / Técnico / Cliente", "Redirección al panel respectivo sin fallos", "Aprobado"],
        ["CP02", "Concurrencia de Horarios", "Dos clientes intentan reservar la misma franja", "El sistema asigna al primero y bloquea al segundo", "Aprobado"],
        ["CP03", "Aviso de Retraso", "Cliente reporta +15 min desde su turno", "La cita pasa a 'Retrasado con Aviso' en el técnico", "Aprobado"],
        ["CP04", "Descarga de Repuesto", "Técnico asigna repuesto a orden finalizada", "Stock disminuye y se genera registro en kardex", "Aprobado"],
        ["CP05", "Auditoría PL/pgSQL", "Actualización o borrado de registro de cita", "Trigger escribe snapshot JSONB en auditoria_log", "Aprobado"]
    ]
    add_table(slide32, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.6), headers32, rows32, [Inches(0.9), Inches(2.5), Inches(3.2), Inches(3.6), Inches(1.5)])
    add_speaker_notes(slide32, "Ejecutamos pruebas funcionales y de concurrencia, validando que el 100% de los casos de prueba superaron los criterios de aceptación.")

    # ==========================================
    # SLIDE 33: Conclusiones e Impacto
    # ==========================================
    slide33 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide33, C_LIGHT)
    add_header(slide33, "Conclusiones y Trabajo Futuro", "Sección IV: Desarrollo")

    add_card(slide33, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Conclusiones del Proyecto", [
        "Se automatizó exitosamente el agendamiento y seguimiento técnico en el taller.",
        "Se eliminó el error humano en inventarios vinculando repuestos a órdenes de trabajo.",
        "La auditoría por triggers en PostgreSQL garantiza total integridad de los datos.",
        "Se cumplieron todas las competencias técnicas exigidas por el programa SENA."
    ], "🎯")

    add_card(slide33, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Trabajo Futuro (Post-MVP)", [
        "Integración de pasarela de pagos electrónicos (Wompi / PayU).",
        "Notificaciones automáticas vía WhatsApp Business API.",
        "Módulo de firma digital del cliente en la recepción de equipos.",
        "Despliegue escalable en contenedores Docker / AWS."
    ], "🚀")

    add_speaker_notes(slide33, "Servitech-app demuestra ser una solución completa, rentable y técnicamente sólida, con una hoja de ruta clara para su evolución futura.")

    # ==========================================
    # SLIDE 34: Cierre / Gracias
    # ==========================================
    slide34 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide34, C_NAVY)

    accent34 = slide34.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.5), Inches(7.5))
    accent34.fill.solid()
    accent34.fill.fore_color.rgb = C_GREEN
    accent34.line.fill.background()

    tb34 = slide34.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.5), Inches(3.5))
    tf34 = tb34.text_frame
    tf34.word_wrap = True

    p = tf34.paragraphs[0]
    p.text = "¡MUCHAS GRACIAS!"
    p.font.name = "Arial"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = C_GREEN
    p.space_after = Pt(14)

    p2 = tf34.add_paragraph()
    p2.text = "Servitech-app • Sistema de Agendamiento y Gestión de Servicios Técnicos"
    p2.font.name = "Arial"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = C_WHITE
    p2.space_after = Pt(10)

    p3 = tf34.add_paragraph()
    p3.text = "Aprendices: Juan Bayona • Cristian Contreras • Emmanuel Flores • Breyner Peña\nSENA Regional — Septiembre 2026\n\n¿Preguntas o comentarios del comité evaluador?"
    p3.font.name = "Arial"
    p3.font.size = Pt(13)
    p3.font.color.rgb = RGBColor(203, 213, 225)

    add_speaker_notes(slide34, "Agradecemos a los instructores por su atención y abrimos el espacio para la sesión de preguntas y retroalimentación.")

    # Save presentation
    output_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "Servitech_Presentacion_SENA.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully at: {output_path}")

if __name__ == "__main__":
    create_deck()
