import os
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def apply_clean_analysis_slides():
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    pptx_path = os.path.join(base_dir, "Servitech-app.pptx")

    prs = pptx.Presentation(pptx_path)

    # ── Paleta de Colores Funsamez / Estándar Limpio ─────────────────────────
    C_DARK = RGBColor(22, 22, 30)           # Texto oscuro (#16161E)
    C_WHITE = RGBColor(255, 255, 255)       # Blanco
    C_RED_CORAL = RGBColor(255, 128, 128)   # Rojo/Coral para Prioridad Alta (25)
    C_ORANGE = RGBColor(255, 194, 102)      # Naranja para Prioridad Media-Alta (16-20)
    C_YELLOW = RGBColor(255, 242, 163)      # Amarillo para Prioridad Media (8-15)
    C_SENA_NAVY = RGBColor(0, 50, 77)       # Azul Marino SENA (#00324D)
    C_LIGHT_BG = RGBColor(248, 250, 252)    # Fondo alternado sutil (#F8FAFC)

    def set_slide_title(slide, title_text):
        title_shape = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape.top.inches < 1.6 and shape.left.inches < 2.0:
                title_shape = shape
                break

        if not title_shape:
            title_shape = slide.shapes.add_textbox(Inches(0.92), Inches(0.40), Inches(11.5), Inches(1.45))

        tf = title_shape.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Arial"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = C_DARK

    def clean_slide_body(slide):
        """Elimina formas del cuerpo conservando el título y el logo superior derecho."""
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.left.inches >= 11.5 and shape.top.inches <= 1.0:
                continue
            if shape.has_text_frame and shape.top.inches < 1.4 and shape.left.inches < 2.0:
                continue
            shapes_to_remove.append(shape)

        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)

    def set_speaker_notes(slide, text):
        slide.notes_slide.notes_text_frame.text = text

    # =========================================================================
    # SLIDE 12: Portada FASE II / ANÁLISIS
    # =========================================================================
    s12 = prs.slides[11]
    set_speaker_notes(s12, "Fase II: Análisis de requerimientos del sistema ServiTech, estructurado bajo el modelo de negocio del taller técnico, fuentes de elicitación, matriz de stakeholders, historias de usuario, requerimientos funcionales, no funcionales y matriz de priorización.")

    # =========================================================================
    # SLIDE 13: Ingeniería de Requisitos (2 COLUMNAS NATIVAS LIMPIAS)
    # =========================================================================
    s13 = prs.slides[12]
    set_slide_title(s13, "Ingeniería de Requisitos")
    clean_slide_body(s13)

    # Columna Izquierda: Elicitación y Fuentes
    box_l = s13.shapes.add_textbox(Inches(0.92), Inches(1.85), Inches(5.60), Inches(5.00))
    tf_l = box_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_top = tf_l.margin_right = tf_l.margin_bottom = 0
    tf_l.clear()

    p_l0 = tf_l.paragraphs[0]
    p_l0.text = "Elicitación de requisitos"
    p_l0.font.name = "Arial"
    p_l0.font.size = Pt(20)
    p_l0.font.bold = True
    p_l0.font.color.rgb = C_DARK
    p_l0.space_after = Pt(6)

    p_l1 = tf_l.add_paragraph()
    p_l1.text = "Fuentes identificadas"
    p_l1.font.name = "Arial"
    p_l1.font.size = Pt(16)
    p_l1.font.bold = True
    p_l1.font.color.rgb = C_DARK
    p_l1.space_after = Pt(12)

    fuentes = [
        "• Personal técnico especializado del taller (Niveles N1, N2, N3).",
        "• Clientes particulares y empresas de servicio técnico.",
        "• Gerencia y propietarios del taller técnico.",
        "• Distribuidores y proveedores de repuestos y componentes."
    ]
    for f in fuentes:
        p = tf_l.add_paragraph()
        p.text = f
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.color.rgb = C_DARK
        p.space_after = Pt(8)

    # Columna Derecha: Técnicas e Instrumentos
    box_r = s13.shapes.add_textbox(Inches(6.82), Inches(1.85), Inches(5.60), Inches(5.00))
    tf_r = box_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_top = tf_r.margin_right = tf_r.margin_bottom = 0
    tf_r.clear()

    p_r0 = tf_r.paragraphs[0]
    p_r0.text = "Técnicas e instrumentos para elicitar"
    p_r0.font.name = "Arial"
    p_r0.font.size = Pt(20)
    p_r0.font.bold = True
    p_r0.font.color.rgb = C_DARK
    p_r0.space_after = Pt(18)

    tecnicas = [
        "• 3 encuestas: información general, historias de usuario y priorización de requisitos.",
        "• 2 entrevistas con administradores y técnicos de boxes (presencial y diagnóstico de SLA).",
        "• Observación directa del flujo de trabajo: recepción -> diagnóstico -> repuestos -> entrega.",
        "• Digitalización y normalización de boletas de garantía físicas y cuadernos de turnos."
    ]
    for t in tecnicas:
        p = tf_r.add_paragraph()
        p.text = t
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.color.rgb = C_DARK
        p.space_after = Pt(8)

    set_speaker_notes(s13, "Elicitación de requisitos basada en fuentes reales del taller y aplicación de encuestas, entrevistas y observación directa para definir las reglas de negocio.")

    # =========================================================================
    # SLIDE 14: Matriz Stakeholders (TEXTO NATIVO LIMPIO SIN FORMAS ARTIFICIALES)
    # =========================================================================
    s14 = prs.slides[13]
    set_slide_title(s14, "Matriz Stakeholders")
    clean_slide_body(s14)

    box14 = s14.shapes.add_textbox(Inches(0.92), Inches(1.85), Inches(11.50), Inches(5.00))
    tf14 = box14.text_frame
    tf14.word_wrap = True
    tf14.margin_left = tf14.margin_top = tf14.margin_right = tf14.margin_bottom = 0
    tf14.clear()

    stakeholders_list = [
        ("1. Satisfacer (Alto Poder / Bajo Interés):",
         "Proveedores de Repuestos y Componentes — Distribuidores comerciales de pantallas y repuestos con acuerdos de SLA y gestión de garantías."),
        
        ("2. Gestionar atentamente (Alto Poder / Alto Interés):",
         "Gerente / Dueño del Taller (Administrador General) — Toma de decisiones estratégicas, rentabilidad, tarifas SLA, control total de mermas y auditoría contable."),
        
        ("3. Informar / Comunicar (Bajo Poder / Alto Interés):",
         "Jefe de Taller y Técnicos (N1, N2, N3) — Operación en boxes de trabajo, recepción de visitas presenciales (walk-ins) y cumplimiento de tiempos de entrega.\n• Clientes Particulares y Empresas — Agendamiento autónomo 24/7, trazabilidad en vivo de tickets y notificación de retrasos."),
        
        ("4. Monitorear (Bajo Poder / Bajo Interés):",
         "Visitantes Web y Usuarios Ocasionales — Consulta pública de tarifas de catálogo y ubicación del taller sin registro activo.")
    ]

    for idx, (title_stk, desc_stk) in enumerate(stakeholders_list):
        p = tf14.add_paragraph() if idx > 0 else tf14.paragraphs[0]
        p.space_after = Pt(10)
        
        r1 = p.add_run()
        r1.text = f"• {title_stk} "
        r1.font.name = "Arial"
        r1.font.size = Pt(15)
        r1.font.bold = True
        r1.font.color.rgb = C_DARK
        
        r2 = p.add_run()
        r2.text = desc_stk
        r2.font.name = "Arial"
        r2.font.size = Pt(14)
        r2.font.color.rgb = C_DARK

    set_speaker_notes(s14, "Matriz de stakeholders basada en el modelo de poder e interés identificando los 4 cuadrantes de gestión del taller técnico.")

    # =========================================================================
    # SLIDE 15: Herramienta Captura 1 (HU01 - TEXTO LIMPIO Y GRANDE)
    # =========================================================================
    s15 = prs.slides[14]
    set_slide_title(s15, "Herramienta para la Captura de Requisitos")
    clean_slide_body(s15)

    box15 = s15.shapes.add_textbox(Inches(1.20), Inches(1.85), Inches(10.80), Inches(5.00))
    tf15 = box15.text_frame
    tf15.word_wrap = True
    tf15.margin_left = tf15.margin_top = tf15.margin_right = tf15.margin_bottom = 0
    tf15.clear()

    hu01_items = [
        ("HU01 - Agendamiento de servicio técnico.", "", 22, True, 8),
        ("Como:", "cliente.", 16, True, 4),
        ("Quiero:", "seleccionar mi tipo de dispositivo (Celular / Laptop / PC), el servicio requerido del catálogo y reservar una franja horaria disponible.", 15, True, 4),
        ("Para:", "asegurar la atención puntual de mi equipo en el taller sin hacer filas presenciales y conocer el tiempo estimado de entrega.", 15, True, 10),
        ("Condiciones:", "", 16, True, 4),
        ("• El sistema debe mostrar el catálogo de servicios con duraciones estimadas (45 a 90 min) y costo base.", "", 14, False, 3),
        ("• El sistema debe validar en tiempo real la disponibilidad de la agenda del técnico evitando colisiones de turnos.", "", 14, False, 3),
        ("• El sistema debe solicitar obligatoriamente marca, modelo y falla reportada del equipo antes de confirmar.", "", 14, False, 3),
        ("• El sistema debe generar confirmación inmediata en pantalla con el número de cita y resumen del servicio.", "", 14, False, 8),
        ("Prioridad:", "25", 16, True, 0)
    ]

    for idx, (h, desc, sz, bld, sp) in enumerate(hu01_items):
        p = tf15.add_paragraph() if idx > 0 else tf15.paragraphs[0]
        p.space_after = Pt(sp)
        if h:
            r1 = p.add_run()
            r1.text = f"{h} " if desc else h
            r1.font.name = "Arial"
            r1.font.size = Pt(sz)
            r1.font.bold = bld
            r1.font.color.rgb = C_DARK
        if desc:
            r2 = p.add_run()
            r2.text = desc
            r2.font.name = "Arial"
            r2.font.size = Pt(sz)
            r2.font.bold = False
            r2.font.color.rgb = C_DARK

    set_speaker_notes(s15, "Formato formal de Historia de Usuario HU01 que captura el agendamiento autónomo del cliente con sus criterios de aceptación.")

    # =========================================================================
    # SLIDE 16: Herramienta Captura 2 (HU03 - TEXTO LIMPIO Y GRANDE)
    # =========================================================================
    s16 = prs.slides[15]
    set_slide_title(s16, "Herramienta para la Captura de Requisitos")
    clean_slide_body(s16)

    box16 = s16.shapes.add_textbox(Inches(1.20), Inches(1.85), Inches(10.80), Inches(5.00))
    tf16 = box16.text_frame
    tf16.word_wrap = True
    tf16.margin_left = tf16.margin_top = tf16.margin_right = tf16.margin_bottom = 0
    tf16.clear()

    hu03_items = [
        ("HU03 - Notificación de retraso en tiempo real.", "", 22, True, 8),
        ("Como:", "cliente con una cita agendada en el taller.", 16, True, 4),
        ("Quiero:", "presionar el botón 'Llegaré tarde (+10 / +15 min)' desde el enlace de recordatorio de mi turno.", 15, True, 4),
        ("Para:", "avisar al taller sobre mi demora y evitar que cancelen o liberen mi cupo de atención.", 15, True, 10),
        ("Condiciones:", "", 16, True, 4),
        ("• La acción debe realizarse mediante un enlace seguro de acceso directo al turno sin necesidad de login complejo.", "", 14, False, 3),
        ("• El estado de la cita debe actualizarse automáticamente a 'Retrasado con Aviso' en el panel del técnico.", "", 14, False, 3),
        ("• El sistema debe otorgar un margen de tolerancia de 15 minutos en la mesa de trabajo antes de marcar inasistencia.", "", 14, False, 3),
        ("• En caso de cancelación por parte del cliente, el sistema debe liberar inmediatamente el cupo para clientes presenciales (walk-ins).", "", 14, False, 8),
        ("Prioridad:", "16", 16, True, 0)
    ]

    for idx, (h, desc, sz, bld, sp) in enumerate(hu03_items):
        p = tf16.add_paragraph() if idx > 0 else tf16.paragraphs[0]
        p.space_after = Pt(sp)
        if h:
            r1 = p.add_run()
            r1.text = f"{h} " if desc else h
            r1.font.name = "Arial"
            r1.font.size = Pt(sz)
            r1.font.bold = bld
            r1.font.color.rgb = C_DARK
        if desc:
            r2 = p.add_run()
            r2.text = desc
            r2.font.name = "Arial"
            r2.font.size = Pt(sz)
            r2.font.bold = False
            r2.font.color.rgb = C_DARK

    set_speaker_notes(s16, "Formato formal de Historia de Usuario HU03 que gestiona el aviso de retraso y las contingencias operativas del taller.")

    # =========================================================================
    # SLIDE 17: Requerimientos Funcionales (TABLA PURA ESTILO EJEMPLO)
    # =========================================================================
    s17 = prs.slides[16]
    set_slide_title(s17, "Requerimientos Funcionales")
    clean_slide_body(s17)

    headers17 = ["TIPO DE REQUISITO", "DESCRIPCIÓN", "CRITERIO DE ACEPTACIÓN", "PRIORIDAD"]
    row17_data = [
        "HU01-RF-05\nFUNCIONAL",
        "El sistema debe registrar en el motor de agendamiento cada cita de servicio técnico en tiempo real.",
        "Cada cita debe registrarse automáticamente en la base de datos relacional.\nLa información registrada debe incluir:\n• ID único de la cita.\n• Nombre completo del cliente solicitante.\n• Dispositivo asociado (marca, modelo, serial/IMEI).\n• Servicio seleccionado y duración SLA estimada (45-90 min).\n• Técnico asignado (o modo 'Cualquier técnico').\n• Estado inicial de la orden (Confirmada).",
        "25"
    ]
    col_w17 = [Inches(2.20), Inches(3.80), Inches(4.50), Inches(1.00)]

    t_shape17 = s17.shapes.add_table(2, 4, Inches(0.92), Inches(2.20), Inches(11.50), Inches(4.30))
    t17 = t_shape17.table
    for i, w in enumerate(col_w17):
        t17.columns[i].width = w

    for c_idx, h_text in enumerate(headers17):
        cell = t17.cell(0, c_idx)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_WHITE
        cell.margin_left = cell.margin_right = Inches(0.10)
        cell.margin_top = cell.margin_bottom = Inches(0.06)
        p = cell.text_frame.paragraphs[0]
        p.text = h_text
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = C_DARK
        p.alignment = PP_ALIGN.CENTER

    for c_idx, val in enumerate(row17_data):
        cell = t17.cell(1, c_idx)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        if c_idx == 3:
            cell.fill.fore_color.rgb = C_RED_CORAL  # Fondo Coral para prioridad 25 igual a Slide 19
        else:
            cell.fill.fore_color.rgb = C_WHITE

        cell.margin_left = cell.margin_right = Inches(0.10)
        cell.margin_top = cell.margin_bottom = Inches(0.06)
        tf = cell.text_frame
        tf.word_wrap = True
        tf.clear()
        for l_idx, line in enumerate(val.split("\n")):
            p = tf.add_paragraph() if l_idx > 0 else tf.paragraphs[0]
            p.text = line
            p.font.name = "Arial"
            p.font.size = Pt(12)
            p.font.color.rgb = C_DARK
            if c_idx in [0, 3]:
                p.alignment = PP_ALIGN.CENTER
                p.font.bold = True

    set_speaker_notes(s17, "Especificación formal en tabla del requerimiento funcional principal RF-05 del motor de agendamiento.")

    # =========================================================================
    # SLIDE 18: Requerimientos no Funcionales (TABLA PURA ESTILO EJEMPLO)
    # =========================================================================
    s18 = prs.slides[17]
    set_slide_title(s18, "Requerimientos no Funcionales")
    clean_slide_body(s18)

    sub_box18 = s18.shapes.add_textbox(Inches(0.92), Inches(1.80), Inches(11.50), Inches(0.45))
    p_sub18 = sub_box18.text_frame.paragraphs[0]
    p_sub18.text = "Especificación de Requisito no Funcional RNF-03"
    p_sub18.font.name = "Arial"
    p_sub18.font.size = Pt(16)
    p_sub18.font.bold = True
    p_sub18.font.color.rgb = C_DARK

    headers18 = ["ID", "TIPO DE REQUISITO", "DESCRIPCIÓN", "CRITERIO DE ACEPTACIÓN", "PRIORIDAD"]
    row18_data = [
        "RNF-03", "NO FUNCIONAL",
        "El sistema debe garantizar la integridad transaccional y la auditoría automática en base de datos.",
        "El sistema debe implementar triggers nativos PL/pgSQL en PostgreSQL que intercepten operaciones INSERT, UPDATE y DELETE.\nCada registro debe almacenarse de forma inmutable en auditoria_log e incluir:\n• Tabla afectada e ID del registro.\n• Datos anteriores y nuevos en formato JSONB.\n• Usuario de base de datos y dirección IP.\n• Timestamp de la operación.\n• Sobrecarga de ejecución inferior a 5 ms.",
        "Alta"
    ]
    col_w18 = [Inches(1.20), Inches(1.80), Inches(3.30), Inches(4.20), Inches(1.00)]

    t_shape18 = s18.shapes.add_table(2, 5, Inches(0.92), Inches(2.40), Inches(11.50), Inches(4.20))
    t18 = t_shape18.table
    for i, w in enumerate(col_w18):
        t18.columns[i].width = w

    for c_idx, h_text in enumerate(headers18):
        cell = t18.cell(0, c_idx)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_WHITE
        cell.margin_left = cell.margin_right = Inches(0.10)
        cell.margin_top = cell.margin_bottom = Inches(0.06)
        p = cell.text_frame.paragraphs[0]
        p.text = h_text
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = C_DARK
        p.alignment = PP_ALIGN.CENTER

    for c_idx, val in enumerate(row18_data):
        cell = t18.cell(1, c_idx)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        if c_idx == 4:
            cell.fill.fore_color.rgb = C_RED_CORAL  # Fondo Coral para prioridad Alta
        else:
            cell.fill.fore_color.rgb = C_WHITE

        cell.margin_left = cell.margin_right = Inches(0.10)
        cell.margin_top = cell.margin_bottom = Inches(0.06)
        tf = cell.text_frame
        tf.word_wrap = True
        tf.clear()
        for l_idx, line in enumerate(val.split("\n")):
            p = tf.add_paragraph() if l_idx > 0 else tf.paragraphs[0]
            p.text = line
            p.font.name = "Arial"
            p.font.size = Pt(12)
            p.font.color.rgb = C_DARK
            if c_idx in [0, 1, 4]:
                p.alignment = PP_ALIGN.CENTER
                p.font.bold = True

    set_speaker_notes(s18, "Especificación formal en tabla del requerimiento no funcional RNF-03 bajo norma ISO 25010 para auditoría automática por triggers PostgreSQL.")

    # =========================================================================
    # SLIDE 19: Priorización MoSCoW (TABLA EXACTA CON COLUMNA SECTOR COLOREADA)
    # =========================================================================
    s19 = prs.slides[18]
    set_slide_title(s19, "Técnica de Priorización de Requisitos")
    clean_slide_body(s19)

    sub_box19 = s19.shapes.add_textbox(Inches(0.92), Inches(1.80), Inches(11.50), Inches(0.50))
    tf_s19 = sub_box19.text_frame
    tf_s19.word_wrap = True
    tf_s19.margin_left = tf_s19.margin_top = tf_s19.margin_right = tf_s19.margin_bottom = 0
    tf_s19.clear()
    p_s19 = tf_s19.paragraphs[0]
    p_s19.text = "Resultado de la Priorización para las Historias de Usuario"
    p_s19.font.name = "Arial"
    p_s19.font.size = Pt(18)
    p_s19.font.bold = False
    p_s19.font.color.rgb = C_DARK

    headers19 = ["Requerimientos", "Valor de negocio", "Urgencia", "Sector"]
    rows19_data = [
        ("HU01", 5, 5, 25, C_RED_CORAL),
        ("HU02", 5, 4, 20, C_ORANGE),
        ("HU03", 4, 4, 16, C_ORANGE),
        ("HU04", 4, 3, 12, C_YELLOW),
        ("HU05", 5, 5, 25, C_RED_CORAL),
        ("HU06", 5, 4, 20, C_ORANGE),
        ("HU07", 3, 3, 9,  C_YELLOW),
        ("HU08", 5, 4, 20, C_ORANGE)
    ]

    t_shape19 = s19.shapes.add_table(len(rows19_data) + 1, 4, Inches(0.92), Inches(2.60), Inches(11.50), Inches(3.80))
    t19 = t_shape19.table
    col_w19 = [Inches(2.875), Inches(2.875), Inches(2.875), Inches(2.875)]
    for idx, w in enumerate(col_w19):
        t19.columns[idx].width = w

    for col_idx, h_text in enumerate(headers19):
        cell = t19.cell(0, col_idx)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_WHITE
        cell.margin_left = cell.margin_right = Inches(0.08)
        cell.margin_top = cell.margin_bottom = Inches(0.06)
        p = cell.text_frame.paragraphs[0]
        p.text = h_text
        p.font.name = "Arial"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = C_DARK
        p.alignment = PP_ALIGN.CENTER

    for row_idx, (req, val_neg, urg, sec, sec_color) in enumerate(rows19_data):
        row_vals = [req, str(val_neg), str(urg), str(sec)]
        for col_idx, val_str in enumerate(row_vals):
            cell = t19.cell(row_idx + 1, col_idx)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if col_idx == 3:
                cell.fill.fore_color.rgb = sec_color
            else:
                cell.fill.fore_color.rgb = C_WHITE

            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            p.text = val_str
            p.font.name = "Arial"
            p.font.size = Pt(12)
            p.font.color.rgb = C_DARK
            p.alignment = PP_ALIGN.CENTER

    set_speaker_notes(s19, "Presentamos la matriz de priorización ponderada de Historias de Usuario: calculamos el Sector multiplicando Valor de Negocio por Urgencia (escala 1 a 5), clasificando en rojo coral las de máxima prioridad (25) como HU01 y HU05, en naranja las prioritarias (16-20) y en amarillo las complementarias.")

    # =========================================================================
    # FASE III: CASOS DE USO (SLIDES 22 A 25 - TABLAS LIMPIAS FORMATO OFICIAL)
    # =========================================================================
    def add_use_case_table_clean(slide, cu_id, cu_name, actor, precond, desc, normal_flow, alt_flow, postcond):
        clean_slide_body(slide)
        headers = ["Campo", "Detalle"]
        rows = [
            ["ID-CU", cu_id],
            ["Nombre CU", cu_name],
            ["Descripción", desc],
            ["Actor", actor],
            ["Precondiciones", precond],
            ["Flujo normal", normal_flow],
            ["Postcondiciones", postcond],
            ["Flujos alternativos", alt_flow]
        ]
        
        table_shape = slide.shapes.add_table(len(rows) + 1, 2, Inches(0.92), Inches(1.75), Inches(11.50), Inches(5.15))
        table = table_shape.table
        table.columns[0].width = Inches(2.20)
        table.columns[1].width = Inches(9.30)

        # Header Row
        for col_idx, h_text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_SENA_NAVY
            cell.margin_left = cell.margin_right = Inches(0.12)
            cell.margin_top = cell.margin_bottom = Inches(0.06)
            p = cell.text_frame.paragraphs[0]
            p.text = h_text
            p.font.name = "Arial"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = C_WHITE
            if col_idx == 0:
                p.alignment = PP_ALIGN.CENTER

        # Data Rows
        for row_idx, r_data in enumerate(rows):
            bg_color = C_WHITE if row_idx % 2 == 0 else C_LIGHT_BG
            for col_idx, val in enumerate(r_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_color
                cell.margin_left = cell.margin_right = Inches(0.12)
                cell.margin_top = cell.margin_bottom = Inches(0.04)
                
                tf = cell.text_frame
                tf.word_wrap = True
                tf.clear()
                
                for l_idx, line in enumerate(val.split("\n")):
                    p = tf.add_paragraph() if l_idx > 0 else tf.paragraphs[0]
                    p.text = line
                    p.font.name = "Arial"
                    p.font.size = Pt(11.5)
                    p.font.color.rgb = C_DARK
                    if col_idx == 0:
                        p.font.bold = True
                        p.alignment = PP_ALIGN.CENTER

    # SLIDE 22: CU02 - Agendar Cita
    s22 = prs.slides[21]
    set_slide_title(s22, "Casos de Usos tablas")
    add_use_case_table_clean(
        s22,
        cu_id="CU02",
        cu_name="Agendar Cita de Servicio Técnico",
        actor="Cliente",
        precond="El servicio deseado debe estar activo en el catálogo del sistema.",
        desc="Proceso mediante el cual un cliente selecciona el tipo de equipo, servicio requerido y reserva un horario disponible.",
        normal_flow="1. El cliente selecciona la opción 'Nueva Cita'.\n2. El cliente selecciona el tipo de dispositivo a reparar o revisar.\n3. El cliente elige el servicio técnico del catálogo.\n4. El sistema despliega las fechas y franjas horarias disponibles.\n5. El cliente selecciona fecha, hora y confirma la reserva.\n6. El sistema guarda la cita y genera una confirmación en pantalla.",
        alt_flow="5a. La franja horaria es reservada simultáneamente por otro cliente: El sistema notifica la indisponibilidad y solicita elegir otro horario.",
        postcond="La cita queda registrada en estado 'Pendiente' y asignada al flujo de atención."
    )
    set_speaker_notes(s22, "Caso de Uso CU02: Agendamiento de Cita de Servicio Técnico.")

    # SLIDE 23: CU03 - Administrar Estado y Demoras
    s23 = prs.slides[22]
    set_slide_title(s23, "Casos de Usos tablas")
    add_use_case_table_clean(
        s23,
        cu_id="CU03",
        cu_name="Administrar Estado y Demoras de Citas",
        actor="Técnico / Administrador",
        precond="Debe existir al menos una cita registrada en el sistema.",
        desc="Permite gestionar el estado de las citas (confirmar, reprogramar, cancelar) y notificar imprevistos o demoras al cliente.",
        normal_flow="1. El actor ingresa al módulo 'Citas' o 'Mi Agenda'.\n2. El sistema muestra el listado de citas registradas.\n3. El actor selecciona una cita específica.\n4. El actor actualiza el estado o pulsa la opción 'Notificar Demora'.\n5. El actor ingresa el tiempo estimado de retraso y el motivo.\n6. El sistema actualiza la cita y envía la notificación de contingencia.",
        alt_flow="4a. El cliente no asiste tras 15 min: El técnico pulsa 'Liberar por No-Show' y reasigna el box para atención presencial (walk-ins).",
        postcond="El nuevo estado o aviso de contingencia queda registrado en la orden de servicio."
    )
    set_speaker_notes(s23, "Caso de Uso CU03: Administrar Estado y Demoras de Citas.")

    # SLIDE 24: CU04 - Controlar Estado de Turno Técnico
    s24 = prs.slides[23]
    set_slide_title(s24, "Casos de Usos tablas")
    add_use_case_table_clean(
        s24,
        cu_id="CU04",
        cu_name="Controlar Estado de Turno Técnico",
        actor="Técnico",
        precond="El técnico debe tener la sesión activa en el panel.",
        desc="Permite a los técnicos cambiar su estado operativo (Disponible, En Servicio, Pausar Turno) durante la jornada laboral.",
        normal_flow="1. El técnico visualiza su indicador de estado en la barra lateral o panel.\n2. El técnico hace clic en la acción 'Pausar Turno' o cambiar estado.\n3. El sistema valida que el técnico no tenga servicios asignados en ejecución en ese instante.\n4. El sistema actualiza el indicador a 'En Pausa' / 'No Disponible'.",
        alt_flow="3a. El técnico tiene un servicio activo: El sistema muestra un mensaje requiriendo finalizar la orden actual antes de pausar el turno.",
        postcond="El sistema no asigna nuevas citas automáticas mientras el técnico se encuentre en pausa."
    )
    set_speaker_notes(s24, "Caso de Uso CU04: Controlar Estado de Turno Técnico.")

    # SLIDE 25: CU05 - Gestionar Catálogo de Servicios y SLA
    s25 = prs.slides[24]
    set_slide_title(s25, "Casos de Usos tablas")
    add_use_case_table_clean(
        s25,
        cu_id="CU05",
        cu_name="Gestionar Oferta de Servicios y SLA",
        actor="Administrador",
        precond="Contar con permisos de administración general.",
        desc="Permite registrar, editar precios, tiempos estimados de atención y configurar acuerdos de nivel de servicio (SLA).",
        normal_flow="1. El administrador ingresa al módulo 'Catálogo'.\n2. El sistema despliega la lista de servicios activos e inactivos.\n3. El administrador selecciona 'Nuevo Servicio' o edita uno existente.\n4. El administrador especifica nombre, costo base, tiempo estimado (45-90 min) y margen buffer.\n5. El administrador guarda los cambios.\n6. El sistema actualiza el catálogo global.",
        alt_flow="5a. Datos obligatorios incompletos: El sistema resalta los campos faltantes y no permite guardar hasta corregir.",
        postcond="La oferta de servicios queda disponible e inmediatamente actualizada para las reservas de los clientes."
    )
    set_speaker_notes(s25, "Caso de Uso CU05: Gestionar Oferta de Servicios y SLA.")

    prs.save(pptx_path)
    print("Servitech-app.pptx actualizado limpiamente sin diseño artificial y 100% alineado con Funsamez adaptado a ServiTech!")

if __name__ == "__main__":
    apply_clean_analysis_slides()
