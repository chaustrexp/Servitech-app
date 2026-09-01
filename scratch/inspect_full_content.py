import pptx

def inspect_slides_content(filename, indices):
    prs = pptx.Presentation(filename)
    for idx in indices:
        if idx >= len(prs.slides):
            continue
        slide = prs.slides[idx]
        print(f"\n============================================================")
        print(f"SLIDE {idx+1:02d} ({slide.slide_layout.name})")
        print(f"============================================================")
        for s_idx, shape in enumerate(slide.shapes):
            print(f"  Shape {s_idx}: '{shape.name}' type={shape.shape_type} ({shape.left/914400:.2f}, {shape.top/914400:.2f}, {shape.width/914400:.2f}, {shape.height/914400:.2f})")
            if shape.has_table:
                t = shape.table
                print(f"    [TABLE] {len(t.rows)}x{len(t.columns)}")
                for r_idx in range(len(t.rows)):
                    txt = " | ".join([c.text_frame.text.replace('\n', ' ') for c in t.rows[r_idx].cells])
                    print(f"      R{r_idx}: {txt[:80]}")
            elif shape.has_text_frame:
                for p_idx, p in enumerate(shape.text_frame.paragraphs):
                    t = p.text.strip().replace('\n', ' ')
                    if t:
                        print(f"      P{p_idx}: {t[:90]}")
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            print(f"  [NOTES]: {slide.notes_slide.notes_text_frame.text.strip()[:100]}...")

# Inspect Formulación (3-10), Diseño BD (27-29), Desarrollo (31-32), Implementación (33-34)
inspect_slides_content("Servitech-app.pptx", [3, 4, 5, 6, 7, 8, 9, 10, 27, 28, 29, 31, 33])
