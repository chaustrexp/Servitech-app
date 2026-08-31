import pptx

def inspect_exact_details(filename, slide_indices):
    prs = pptx.Presentation(filename)
    print(f"===============================================================")
    print(f"INSPECTING: {filename}")
    print(f"===============================================================")
    for idx in slide_indices:
        if idx >= len(prs.slides):
            continue
        slide = prs.slides[idx]
        print(f"\n--- SLIDE {idx+1}: {slide.shapes[0].text_frame.text if (slide.shapes and slide.shapes[0].has_text_frame) else 'NO TITLE'} ---")
        for s_idx, shape in enumerate(slide.shapes):
            print(f"  Shape {s_idx}: '{shape.name}' type={shape.shape_type} left={shape.left/914400:.2f}\" top={shape.top/914400:.2f}\" w={shape.width/914400:.2f}\" h={shape.height/914400:.2f}\"")
            if shape.has_table:
                t = shape.table
                print(f"    [TABLE] {len(t.rows)} rows x {len(t.columns)} cols")
                for r_idx, row in enumerate(t.rows):
                    for c_idx, cell in enumerate(row.cells):
                        cell_txt = cell.text_frame.text.replace('\n', ' ')
                        p0 = cell.text_frame.paragraphs[0] if cell.text_frame.paragraphs else None
                        f_name = p0.font.name if (p0 and p0.font) else 'None'
                        f_size = p0.font.size.pt if (p0 and p0.font and p0.font.size) else 'None'
                        f_bold = p0.font.bold if (p0 and p0.font) else 'None'
                        print(f"      [{r_idx},{c_idx}] (Font: {f_name}, {f_size}pt, bold={f_bold}): {cell_txt[:50]}")
            elif shape.has_text_frame:
                for p_idx, p in enumerate(shape.text_frame.paragraphs):
                    p_txt = p.text.strip().replace('\n', ' ')
                    if p_txt:
                        f_name = p.font.name if p.font else 'None'
                        f_size = p.font.size.pt if (p.font and p.font.size) else 'None'
                        f_bold = p.font.bold if p.font else 'None'
                        print(f"      [Para {p_idx}] (Font: {f_name}, {f_size}pt, bold={f_bold}): {p_txt[:60]}")

inspect_exact_details('Servitech-app.pptx', [3, 4, 5, 6, 7, 8, 9, 12, 13, 14, 15, 16, 17, 18, 21, 22])
