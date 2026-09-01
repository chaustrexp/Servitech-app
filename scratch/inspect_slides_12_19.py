import pptx

prs = pptx.Presentation('Servitech-app.pptx')
for idx in range(11, 19):
    slide = prs.slides[idx]
    print(f"\n=== Slide {idx+1} ===")
    for s in slide.shapes:
        print(f"Shape: {s.name}, type={s.shape_type}, left={s.left/914400:.2f}, top={s.top/914400:.2f}, w={s.width/914400:.2f}, h={s.height/914400:.2f}")
        if s.has_text_frame:
            print(f"  Text: {s.text_frame.text}")
