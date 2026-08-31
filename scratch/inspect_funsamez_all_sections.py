import pptx

prs_fun = pptx.Presentation('funsamez_sustentacion_final (2).pptx')

for idx in [3, 4, 5, 6, 7, 8, 9, 10, 25, 26, 27, 28, 29, 30, 31, 32]:
    if idx >= len(prs_fun.slides):
        continue
    slide = prs_fun.slides[idx]
    print(f"\n============================================================")
    print(f"FUNSAMEZ SLIDE {idx+1:02d} ({slide.slide_layout.name})")
    print(f"============================================================")
    for s_idx, shape in enumerate(slide.shapes):
        print(f"  Shape {s_idx}: '{shape.name}' type={shape.shape_type} ({shape.left/914400:.2f}, {shape.top/914400:.2f}, {shape.width/914400:.2f}, {shape.height/914400:.2f})")
        if shape.has_table:
            t = shape.table
            print(f"    [TABLE] {len(t.rows)}x{len(t.columns)}")
            for r_idx in range(len(t.rows)):
                txt = " | ".join([c.text_frame.text.replace('\n', ' ') for c in t.rows[r_idx].cells])
                print(f"      R{r_idx}: {txt[:80]}")
        elif shape.has_text_frame:
            for p_idx, p in enumerate(shape.text_frame.paragraphs):
                t = p.text.strip().replace('\n', ' ')
                if t:
                    print(f"      P{p_idx}: {t}")
