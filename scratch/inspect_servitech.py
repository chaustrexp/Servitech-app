import pptx

prs = pptx.Presentation('Servitech-app.pptx')
print(f'Slide dimensions: width={prs.slide_width.inches:.2f} in, height={prs.slide_height.inches:.2f} in')

for i in range(len(prs.slides)):
    slide = prs.slides[i]
    print(f'\n--- Slide {i+1} ---')
    for s in slide.shapes:
        print(f'Shape {s.name} (type {s.shape_type}): left={s.left/914400:.2f} in, top={s.top/914400:.2f} in, w={s.width/914400:.2f} in, h={s.height/914400:.2f} in')
        if s.has_text_frame:
            for p in s.text_frame.paragraphs:
                font_name = p.font.name if p.font else 'None'
                font_size = p.font.size.pt if (p.font and p.font.size) else 'None'
                txt = p.text.strip()
                if txt:
                    print(f'   [Font: {font_name}, size: {font_size}] {txt[:80]}')
