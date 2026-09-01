import pptx

prs = pptx.Presentation('Servitech-app.pptx')
for idx in range(3, 11):
    slide = prs.slides[idx]
    print(f"\n=== Slide {idx+1} ({slide.slide_layout.name}) ===")
    for s_idx, shape in enumerate(slide.shapes):
        print(f"  Shape {s_idx}: name='{shape.name}', type={shape.shape_type}, pos=({shape.left/914400:.2f}\", {shape.top/914400:.2f}\", {shape.width/914400:.2f}\", {shape.height/914400:.2f}\")")
        if shape.has_text_frame:
            for p_idx, p in enumerate(shape.text_frame.paragraphs):
                runs = [(r.text, r.font.name, r.font.size.pt if r.font.size else None, r.font.bold, r.font.color.rgb if r.font.color and r.font.color.type==1 else None) for r in p.runs]
                print(f"    P{p_idx} (align={p.alignment}): '{p.text}' | Runs: {runs}")
