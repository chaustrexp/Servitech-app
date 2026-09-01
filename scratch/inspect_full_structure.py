import pptx

def inspect_all_slides(filename):
    prs = pptx.Presentation(filename)
    print(f"============================================================")
    print(f"FILE: {filename} (Total Slides: {len(prs.slides)})")
    print(f"============================================================")
    for idx, slide in enumerate(prs.slides):
        title = "NO_TITLE"
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                title = shape.text_frame.text.strip().replace('\n', ' ')[:60]
                break
        has_table = any(s.has_table for s in slide.shapes)
        has_image = any(s.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE for s in slide.shapes)
        print(f"Slide {idx+1:02d}: Layout='{slide.slide_layout.name}' | Table={has_table} | Img={has_image} | Title: '{title}'")

print("\n--- SERVITECH CURRENT SLIDES ---")
inspect_all_slides("Servitech-app.pptx")

print("\n--- FUNSAMEZ REFERENCE SLIDES ---")
inspect_all_slides("funsamez_sustentacion_final (2).pptx")
