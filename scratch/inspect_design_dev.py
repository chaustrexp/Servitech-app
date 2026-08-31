import pptx

def inspect_design_dev_impl(filename):
    prs = pptx.Presentation(filename)
    indices = [20, 21, 25, 26, 27, 28, 29, 30, 31, 32, 33, 34]
    for idx in indices:
        if idx >= len(prs.slides):
            continue
        slide = prs.slides[idx]
        print(f"\n============================================================")
        print(f"SLIDE {idx+1:02d}: Layout='{slide.slide_layout.name}'")
        print(f"============================================================")
        for s_idx, shape in enumerate(slide.shapes):
            txt = shape.text_frame.text.replace('\n', ' ') if shape.has_text_frame else ''
            print(f"  Shape {s_idx}: '{shape.name}' type={shape.shape_type} ({shape.left/914400:.2f}, {shape.top/914400:.2f}) | Text: {txt[:60]}")
            if shape.has_table:
                t = shape.table
                print(f"    [TABLE] {len(t.rows)}x{len(t.columns)}")
                for r in range(min(4, len(t.rows))):
                    r_txt = " | ".join([c.text_frame.text.replace('\n', ' ') for c in t.rows[r].cells])
                    print(f"      R{r}: {r_txt[:75]}")

inspect_design_dev_impl("Servitech-app.pptx")
