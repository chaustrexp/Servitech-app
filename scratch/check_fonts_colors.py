import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = pptx.Presentation('Servitech-app.pptx')
print(f"Presentation has {len(prs.slides)} slides.")

fonts = set()
colors = set()

for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if p.font.name:
                    fonts.add(p.font.name)
                for r in p.runs:
                    if r.font.name:
                        fonts.add(r.font.name)
                    if r.font.color and r.font.color.type == 1:
                        colors.add(r.font.color.rgb)

print("Fonts found:", fonts)
print("Colors found:", colors)
