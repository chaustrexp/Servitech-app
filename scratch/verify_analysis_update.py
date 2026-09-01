import pptx

prs = pptx.Presentation('Servitech-app.pptx')
print(f"Total slides: {len(prs.slides)}")

for idx in range(11, 19): # slides 12 to 19 (0-indexed 11 to 18)
    slide = prs.slides[idx]
    print(f"\n=======================================================")
    print(f"SLIDE {idx+1}: {slide.slide_layout.name}")
    print(f"=======================================================")
    for s_idx, shape in enumerate(slide.shapes):
        print(f"Shape {s_idx}: '{shape.name}' (type {shape.shape_type})")
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                txt = p.text.strip()
                if txt:
                    print(f"   [Text] {txt[:100]}")
        elif shape.has_table:
            tbl = shape.table
            print(f"   [TABLE {len(tbl.rows)} rows x {len(tbl.columns)} cols]")
            for r_idx, r in enumerate(tbl.rows):
                row_vals = [c.text.strip().replace('\n', ' // ') for c in r.cells]
                print(f"      Row {r_idx}: " + " | ".join(row_vals[:4]))
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
        print(f"   [ORADOR] {slide.notes_slide.notes_text_frame.text.strip()}")
