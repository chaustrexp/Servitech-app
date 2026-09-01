import pptx

prs = pptx.Presentation('Servitech-app.pptx')

print(f"Total slides: {len(prs.slides)}")
for idx in range(11, 19): # slides 12 to 19 (0-indexed 11 to 18)
    slide = prs.slides[idx]
    print(f"\n=== Slide {idx+1} ===")
    print(f"Slide layout: {slide.slide_layout.name}")
    for s_idx, shape in enumerate(slide.shapes):
        print(f"  Shape {s_idx}: name='{shape.name}', type={shape.shape_type}, pos=({shape.left/914400:.2f}\", {shape.top/914400:.2f}\", {shape.width/914400:.2f}\", {shape.height/914400:.2f}\")")
        if shape.has_text_frame:
            print(f"    Text: {shape.text_frame.text.strip()[:100]}")
